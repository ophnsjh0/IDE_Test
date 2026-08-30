import json
import tempfile
from datetime import timedelta
from io import BytesIO, StringIO
from pathlib import Path
from types import SimpleNamespace
from unittest.mock import MagicMock, patch

import anthropic

from django.contrib.auth.models import User
from django.core.files.uploadedfile import SimpleUploadedFile
from django.core.management import call_command
from django.core.management.base import CommandError
from django.test import TestCase, override_settings
from django.utils import timezone

from .models import AppSetting, Case, CaseEmail
from .services import (analyzer, eveng, gmail_sync, help_agent,
                       lab_drivers, lab_probe)
from .services.email_parser import (build_gmail_query, clean_subject,
                                    detect_vendor_and_direction,
                                    extract_device_info, find_ignore_reason,
                                    normalize_body)
from .services.gmail_sync import (SyncInProgress, _find_case,
                                  apply_device_info)


# 지식 추출 목을 위한 모델 이름 — generate_structured_with_model이 (실제 모델, 결과)를
# 돌려주므로, 목도 같은 모양이어야 analyzed_by 기록 경로가 함께 검증된다.
MODEL = 'claude-opus-5'


def make_case(**kwargs):
    defaults = dict(vendor='Arista', status='Open', summary='요약', source='email')
    defaults.update(kwargs)
    return Case.objects.create(**defaults)


def make_email(case, subject, thread_id='thread-x', message_id=None):
    return CaseEmail.objects.create(
        case=case,
        gmail_message_id=message_id or f'msg-{CaseEmail.objects.count() + 1}',
        gmail_thread_id=thread_id,
        direction='outbound',
        sender='eng@ubersys.co.kr',
        recipient='support@arista.com',
        subject=subject,
        subject_ko='',
        body_original='본문',
        body_ko='',
        received_at=timezone.now(),
    )


class FindCaseTests(TestCase):
    """Arista처럼 오픈 메일과 SR 확인 메일이 다른 스레드로 갈릴 때의 매칭."""

    def test_confirmation_mail_matches_original_case_by_embedded_subject(self):
        # 엔지니어 오픈 메일로 생성된 케이스 (SR 번호 없음)
        case = make_case(gmail_thread_id='thread-original')
        make_email(case, '40G Interface Link FLAP', thread_id='thread-original')

        # 벤더 확인 메일: 새 스레드 + SR 번호 + 원본 제목 포함
        found = _find_case(
            '834065', 'thread-sr', 'Arista',
            'New UBER Systems Co. Ltd Case: SR 834065 40G Interface Link FLAP [ ref:!00DA0.!500Kh0 ]',
        )
        self.assertEqual(found, case)

    def test_original_mail_matches_confirmation_case_in_reverse_order(self):
        # SR 확인 메일이 먼저 동기화되어 케이스가 이미 있는 경우
        case = make_case(vendor_case_number='825808', gmail_thread_id='thread-sr')
        make_email(
            case,
            'New UBER Systems Co. Ltd Case: SR 825808 [samsung securities] '
            'Continuous PhyEthtool Logs After EOS Upgrade',
            thread_id='thread-sr',
        )

        found = _find_case(
            None, 'thread-original', 'Arista',
            '[samsung securities] Continuous PhyEthtool Logs After EOS Upgrade',
        )
        self.assertEqual(found, case)

    def test_repeated_notification_subjects_do_not_merge(self):
        # 번호 없는 공지 메일끼리는 제목이 같아도 병합 대상이 아님
        case = make_case(gmail_thread_id='thread-notice-1')
        make_email(case, 'New Field notice email notification', thread_id='thread-notice-1')

        found = _find_case(None, 'thread-notice-2', 'Arista',
                           'New Field notice email notification')
        self.assertIsNone(found)

    def test_short_subject_does_not_match(self):
        case = make_case(gmail_thread_id='thread-short')
        make_email(case, 'Link FLAP', thread_id='thread-short')

        found = _find_case('834065', 'thread-new', 'Arista',
                           'New Case: SR 834065 Link FLAP')
        self.assertIsNone(found)

    def test_other_vendor_is_not_matched(self):
        case = make_case(gmail_thread_id='thread-a10')
        make_email(case, '40G Interface Link FLAP', thread_id='thread-a10')

        found = _find_case('834065', 'thread-new', 'A10',
                           'New UBER Systems Co. Ltd Case: SR 834065 40G Interface Link FLAP')
        self.assertIsNone(found)

    def test_thread_of_merged_case_is_traced_via_email(self):
        # 병합으로 케이스 대표 스레드가 아니게 된 스레드도 이메일로 역추적
        case = make_case(vendor_case_number='834065', gmail_thread_id='thread-sr')
        make_email(case, '40G Interface Link FLAP', thread_id='thread-original')

        found = _find_case(None, 'thread-original', 'Arista', 'Re: 40G Interface Link FLAP')
        self.assertEqual(found, case)

    def test_vendor_case_number_still_matches_first(self):
        case = make_case(vendor_case_number='834065', gmail_thread_id='thread-sr')
        found = _find_case('834065', 'unrelated-thread', 'Arista', '아무 제목')
        self.assertEqual(found, case)


CASE_OPEN_BODY = (
    '1. End customer name: NHN\n'
    '2. Partner/Reseller name: ubersystems\n'
    '3. Hardware Platform: TH1040-F\n'
    '4. Software Version: 6.0.8\n'
    '5. Priority : P2\n'
    '6. Serial Number : TH10154022070160\n'
    '7. Description :\n\n'
    'Hi Team\n'
    'After the following logs occurred, the device failed over.\n'
    'Jul 10 2026 04:00:35 Info [HA]:VRRP-A parid 0 vrid 1 state switch '
    'from 1 to 0 (Standby)\n'
    'Jul 10 2026 04:00:36 Info [HA]:VRRP-A parid 0 vrid 1 received higher '
    'priority advertisement from peer, transitioning to backup state now.\n'
    'Please check the attached show techsupport output and let us know the '
    'root cause of this failover event as soon as possible. Thanks.\n'
)


class FindCaseByBodyTests(TestCase):
    """제목을 바꿔 재발송해 스레드가 갈린 동일 접수 메일의 본문 유사도 매칭."""

    def test_resent_mail_with_new_subject_matches_by_body(self):
        case = make_case(vendor='A10', gmail_thread_id='thread-1')
        make_email(case, '[NHN-6.0.8]Device Failover Occurrence',
                   thread_id='thread-1')
        email = CaseEmail.objects.get(case=case)
        email.body_original = CASE_OPEN_BODY
        email.save()

        found = _find_case(
            None, 'thread-2', 'A10',
            '[NHN-6.0.8][AXMON]:Detected problem in Health Monitor',
            CASE_OPEN_BODY + '\nBest regards',
        )
        self.assertEqual(found, case)
        # 병합 표시가 타임라인에 남는다
        self.assertIn('중복 접수 메일', found.action_steps or '')

    def test_different_issue_creates_new_case(self):
        case = make_case(vendor='A10', gmail_thread_id='thread-1')
        make_email(case, '[NHN-6.0.8]Device Failover Occurrence',
                   thread_id='thread-1')
        email = CaseEmail.objects.get(case=case)
        email.body_original = CASE_OPEN_BODY
        email.save()

        other_body = (
            '1. End customer name: Kakao\n'
            '2. Hardware Platform: TH3350\n'
            '3. Serial Number : TH33500000000001\n'
            'Hello, we observed SNMP polling failures on this device after '
            'enabling the new monitoring profile. The walk stops responding '
            'after roughly ten minutes and only recovers when the agent is '
            'restarted manually. Please advise which debug output you need.\n'
        )
        found = _find_case(None, 'thread-2', 'A10', 'SNMP polling issue', other_body)
        self.assertIsNone(found)

    def test_short_body_is_skipped(self):
        case = make_case(vendor='A10', gmail_thread_id='thread-1')
        make_email(case, '[NHN-6.0.8]Device Failover Occurrence',
                   thread_id='thread-1')
        found = _find_case(None, 'thread-2', 'A10', '제목', '감사합니다.')
        self.assertIsNone(found)

    def test_matching_serial_number_relaxes_threshold(self):
        # 공통부(시리얼 포함) + 서로 다른 꼬리말로 유사도를 약 0.92로 구성
        # (ratio = 공통길이/(공통길이+꼬리길이) 이므로 꼬리를 공통부의 8.7%로)
        common = CASE_OPEN_BODY + 'filler word ' * 30
        tail = int(len(normalize_body(common)) * 0.087)
        body_a = common + 'x' * tail
        body_b = common + 'y' * tail

        case = make_case(vendor='A10', gmail_thread_id='thread-1')
        make_email(case, '[NHN-6.0.8]Device Failover Occurrence',
                   thread_id='thread-1')
        email = CaseEmail.objects.get(case=case)
        email.body_original = body_a
        email.save()

        found = _find_case(None, 'thread-2', 'A10', '재발송 제목', body_b)
        self.assertEqual(found, case)

        # 시리얼이 다르면 완화 없이 0.95가 적용되어 매칭되지 않는다
        email.body_original = body_a.replace('TH10154022070160', 'TH99999999999999')
        email.save()
        case.refresh_from_db()
        found = _find_case(None, 'thread-3', 'A10', '재발송 제목', body_b)
        self.assertIsNone(found)

    def test_other_vendor_body_is_not_compared(self):
        case = make_case(vendor='A10', gmail_thread_id='thread-1')
        make_email(case, '[NHN-6.0.8]Device Failover Occurrence',
                   thread_id='thread-1')
        email = CaseEmail.objects.get(case=case)
        email.body_original = CASE_OPEN_BODY
        email.save()

        found = _find_case(None, 'thread-2', 'Arista', '다른 벤더 재발송', CASE_OPEN_BODY)
        self.assertIsNone(found)


class ExtractDeviceInfoTests(TestCase):
    """메일 본문/제목에서 장비 모델·시리얼·버전 추출."""

    def test_a10_open_template(self):
        info = extract_device_info('[NHN-6.0.8]Device Failover Occurrence', CASE_OPEN_BODY)
        self.assertEqual(info['device_model'], 'TH1040-F')
        self.assertEqual(info['device_serial'], 'TH10154022070160')
        self.assertEqual(info['software_version'], '6.0.8')

    def test_hpe_sn_line_items_joined(self):
        body = (
            'RMA parts list:\n'
            'EC-ADV-AAS-UL, S/N 001BBC04E53A\n'
            'EC-BOOST-AAS-10G, S/N 001BBC04E53B\n'
            'EC-DTD-AAS, S/N 001BBC04E53C\n'
        )
        info = extract_device_info('RMA request', body)
        self.assertEqual(info['device_serial'],
                         '001BBC04E53A, 001BBC04E53B, 001BBC04E53C')

    def test_arista_model_token_and_subject_version(self):
        body = 'We upgraded our DCS-7050SX3-48YC12 switch and see PhyEthtool errors.'
        info = extract_device_info('[samsung-4.32.4M] PhyEthtool errors', body)
        self.assertEqual(info['device_model'], 'DCS-7050SX3-48YC12')
        self.assertEqual(info['software_version'], '4.32.4M')

    def test_no_device_info_returns_empty(self):
        info = extract_device_info('New End of Sale email notification',
                                   'The following products reach end of sale next quarter.')
        self.assertEqual(info, {'device_model': '', 'device_serial': '',
                                'software_version': ''})


class ApplyDeviceInfoTests(TestCase):
    """추출값 반영: 정규식 1차 -> AI 2차, 빈 필드만 채움."""

    def test_regex_first_then_ai_fills_missing(self):
        case = make_case(vendor='A10')
        analysis = {'device_model': 'TH9999', 'device_serial': 'AI-SERIAL',
                    'software_version': '9.9.9'}
        # 본문 정규식에서 모델/시리얼/버전을 모두 찾으므로 AI값은 무시된다
        apply_device_info(case, '[NHN-6.0.8] subject', CASE_OPEN_BODY, analysis)
        self.assertEqual(case.device_model, 'TH1040-F')
        self.assertEqual(case.device_serial, 'TH10154022070160')
        self.assertEqual(case.software_version, '6.0.8')

    def test_ai_value_used_when_regex_misses(self):
        case = make_case(vendor='HPE Aruba')
        analysis = {'device_model': 'Aruba 7205', 'device_serial': '',
                    'software_version': '8.10.0.9'}
        apply_device_info(case, 'Gateway issue', '컨트롤러에서 Role 소실 이슈가 발생했습니다.', analysis)
        self.assertEqual(case.device_model, 'Aruba 7205')
        self.assertEqual(case.software_version, '8.10.0.9')
        self.assertEqual(case.device_serial, '')

    def test_existing_values_are_not_overwritten(self):
        case = make_case(vendor='A10', device_model='TH1040-F')
        apply_device_info(case, 'subject', 'Model: TH3350', None)
        self.assertEqual(case.device_model, 'TH1040-F')


class AuthTests(TestCase):
    """세션 인증: 로그인 없이는 API 접근 불가, 로그인/로그아웃 플로우."""

    def setUp(self):
        from django.contrib.auth.models import User
        User.objects.create_user('eng1', password='test-pass-123!')

    def test_api_requires_login(self):
        response = self.client.get('/api/cases/')
        self.assertIn(response.status_code, (401, 403))

    def test_health_check_is_open(self):
        response = self.client.get('/api/health/')
        self.assertEqual(response.status_code, 200)

    def test_login_grants_access_and_me_reports_user(self):
        response = self.client.post('/api/auth/login/',
                                    {'username': 'eng1', 'password': 'test-pass-123!'},
                                    content_type='application/json')
        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.json()['username'], 'eng1')
        # csrftoken 쿠키가 함께 발급된다
        self.assertIn('csrftoken', response.cookies)

        self.assertEqual(self.client.get('/api/cases/').status_code, 200)
        me = self.client.get('/api/auth/me/').json()
        self.assertTrue(me['authenticated'])

    def test_wrong_password_rejected(self):
        response = self.client.post('/api/auth/login/',
                                    {'username': 'eng1', 'password': 'wrong'},
                                    content_type='application/json')
        self.assertEqual(response.status_code, 401)

    def test_me_reports_anonymous_without_session(self):
        me = self.client.get('/api/auth/me/')
        self.assertEqual(me.status_code, 200)
        self.assertFalse(me.json()['authenticated'])

    def test_logout_revokes_session(self):
        self.client.post('/api/auth/login/',
                         {'username': 'eng1', 'password': 'test-pass-123!'},
                         content_type='application/json')
        self.assertEqual(self.client.post('/api/auth/logout/').status_code, 200)
        self.assertIn(self.client.get('/api/cases/').status_code, (401, 403))


class UserManagementTests(TestCase):
    """관리자 전용 계정 발급/관리 API."""

    def setUp(self):
        from django.contrib.auth.models import User
        User.objects.create_user('staff1', password='admin-pass-123!', is_staff=True)
        User.objects.create_user('normal1', password='normal-pass-123!')

    def login(self, username, password):
        return self.client.post('/api/auth/login/',
                                {'username': username, 'password': password},
                                content_type='application/json')

    def test_normal_user_cannot_access(self):
        self.login('normal1', 'normal-pass-123!')
        self.assertEqual(self.client.get('/api/auth/users/').status_code, 403)
        self.assertEqual(self.client.post('/api/auth/users/', {},
                                          content_type='application/json').status_code, 403)

    def test_admin_creates_account(self):
        self.login('staff1', 'admin-pass-123!')
        response = self.client.post(
            '/api/auth/users/',
            {'username': 'eng2', 'password': 'good-pass-77!', 'name': '김엔지니어'},
            content_type='application/json')
        self.assertEqual(response.status_code, 201)
        self.assertEqual(response.json()['username'], 'eng2')
        self.assertEqual(response.json()['role'], 'viewer')
        # 발급된 계정으로 로그인 가능
        self.client.post('/api/auth/logout/')
        self.assertEqual(self.login('eng2', 'good-pass-77!').status_code, 200)

    def test_duplicate_username_rejected(self):
        self.login('staff1', 'admin-pass-123!')
        response = self.client.post('/api/auth/users/',
                                    {'username': 'Normal1', 'password': 'good-pass-77!'},
                                    content_type='application/json')
        self.assertEqual(response.status_code, 400)

    def test_weak_password_rejected(self):
        self.login('staff1', 'admin-pass-123!')
        response = self.client.post('/api/auth/users/',
                                    {'username': 'eng3', 'password': '1234'},
                                    content_type='application/json')
        self.assertEqual(response.status_code, 400)

    def test_deactivate_blocks_login_and_self_deactivation_denied(self):
        from django.contrib.auth.models import User
        self.login('staff1', 'admin-pass-123!')
        normal = User.objects.get(username='normal1')
        staff = User.objects.get(username='staff1')

        # 자기 자신 비활성화는 거부
        response = self.client.patch(f'/api/auth/users/{staff.id}/',
                                     {'is_active': False}, content_type='application/json')
        self.assertEqual(response.status_code, 400)

        # 다른 계정 비활성화 -> 로그인 차단
        response = self.client.patch(f'/api/auth/users/{normal.id}/',
                                     {'is_active': False}, content_type='application/json')
        self.assertEqual(response.status_code, 200)
        self.client.post('/api/auth/logout/')
        self.assertEqual(self.login('normal1', 'normal-pass-123!').status_code, 401)

    def test_password_reset(self):
        from django.contrib.auth.models import User
        self.login('staff1', 'admin-pass-123!')
        normal = User.objects.get(username='normal1')
        response = self.client.patch(f'/api/auth/users/{normal.id}/',
                                     {'password': 'new-pass-88!'},
                                     content_type='application/json')
        self.assertEqual(response.status_code, 200)
        self.client.post('/api/auth/logout/')
        self.assertEqual(self.login('normal1', 'new-pass-88!').status_code, 200)

    def test_delete_account(self):
        from django.contrib.auth.models import User
        self.login('staff1', 'admin-pass-123!')
        normal = User.objects.get(username='normal1')
        response = self.client.delete(f'/api/auth/users/{normal.id}/')
        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.json()['deleted'], 'normal1')
        self.assertFalse(User.objects.filter(username='normal1').exists())
        # 삭제된 계정은 로그인 불가
        self.client.post('/api/auth/logout/')
        self.assertEqual(self.login('normal1', 'normal-pass-123!').status_code, 401)

    def test_self_deletion_denied(self):
        from django.contrib.auth.models import User
        self.login('staff1', 'admin-pass-123!')
        staff = User.objects.get(username='staff1')
        response = self.client.delete(f'/api/auth/users/{staff.id}/')
        self.assertEqual(response.status_code, 400)
        self.assertTrue(User.objects.filter(username='staff1').exists())

    def test_normal_user_cannot_delete(self):
        from django.contrib.auth.models import User
        self.login('normal1', 'normal-pass-123!')
        staff = User.objects.get(username='staff1')
        response = self.client.delete(f'/api/auth/users/{staff.id}/')
        self.assertEqual(response.status_code, 403)


class RolePermissionTests(TestCase):
    """역할(viewer/engineer/admin)별 API 권한 경계."""

    def setUp(self):
        from django.contrib.auth.models import User
        from .permissions import set_user_role
        for username, role in (('v1', 'viewer'), ('e1', 'engineer'), ('a1', 'admin')):
            user = User.objects.create_user(username, password='role-pass-123!')
            set_user_role(user, role)
        self.case = make_case(vendor='A10', summary='권한 테스트용 케이스')

    def login(self, username):
        self.client.post('/api/auth/login/',
                         {'username': username, 'password': 'role-pass-123!'},
                         content_type='application/json')

    def test_viewer_can_read_but_not_write(self):
        self.login('v1')
        self.assertEqual(self.client.get('/api/cases/').status_code, 200)
        self.assertEqual(self.client.get(f'/api/cases/{self.case.id}/').status_code, 200)
        self.assertEqual(self.client.get('/api/dashboard/stats/').status_code, 200)

        create = self.client.post('/api/cases/', {'vendor': 'A10', 'summary': '뷰어 생성 시도'},
                                  content_type='application/json')
        self.assertEqual(create.status_code, 403)
        patch = self.client.patch(f'/api/cases/{self.case.id}/', {'status': 'Resolved'},
                                  content_type='application/json')
        self.assertEqual(patch.status_code, 403)
        sync = self.client.post('/api/gmail/sync/')
        self.assertEqual(sync.status_code, 403)

    def test_engineer_can_write_but_not_delete_or_configure(self):
        self.login('e1')
        patch = self.client.patch(f'/api/cases/{self.case.id}/', {'status': 'Resolved'},
                                  content_type='application/json')
        self.assertEqual(patch.status_code, 200)

        self.assertEqual(self.client.delete(f'/api/cases/{self.case.id}/').status_code, 403)
        model_put = self.client.put('/api/settings/translation-model/', {'model': 'default'},
                                    content_type='application/json')
        self.assertEqual(model_put.status_code, 403)
        self.assertEqual(self.client.get('/api/auth/users/').status_code, 403)

    def test_admin_can_delete_case(self):
        self.login('a1')
        self.assertEqual(self.client.delete(f'/api/cases/{self.case.id}/').status_code, 204)

    def test_gmail_sync_switch_is_readable_by_all_but_admin_only_writable(self):
        url = '/api/settings/gmail-sync/'
        self.login('v1')
        payload = self.client.get(url).json()
        self.assertTrue(payload['enabled'])          # 미설정이면 켜짐이 기본
        self.assertIn('schedule', payload)
        self.assertEqual(self.client.put(url, {'enabled': False},
                                         content_type='application/json').status_code, 403)

        self.login('e1')
        self.assertEqual(self.client.put(url, {'enabled': False},
                                         content_type='application/json').status_code, 403)

        self.login('a1')
        response = self.client.put(url, {'enabled': False}, content_type='application/json')
        self.assertEqual(response.status_code, 200)
        self.assertFalse(response.json()['enabled'])
        self.assertFalse(gmail_sync.is_cron_enabled())

        # 잘못된 값은 거부하고 기존 상태를 유지
        self.assertEqual(self.client.put(url, {'enabled': 'yes'},
                                         content_type='application/json').status_code, 400)
        self.assertFalse(gmail_sync.is_cron_enabled())

    def test_admin_changes_role_but_cannot_demote_self(self):
        from django.contrib.auth.models import User
        self.login('a1')
        e1 = User.objects.get(username='e1')
        response = self.client.patch(f'/api/auth/users/{e1.id}/', {'role': 'viewer'},
                                     content_type='application/json')
        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.json()['role'], 'viewer')

        a1 = User.objects.get(username='a1')
        response = self.client.patch(f'/api/auth/users/{a1.id}/', {'role': 'engineer'},
                                     content_type='application/json')
        self.assertEqual(response.status_code, 400)

    def test_me_includes_role(self):
        self.login('e1')
        me = self.client.get('/api/auth/me/').json()
        self.assertEqual(me['role'], 'engineer')
        self.assertFalse(me['is_admin'])


class ChatSessionTests(TestCase):
    """AI 도우미 대화 저장(ChatSession/ChatTurn)과 엔지니어 이상 권한."""

    FAKE_RESULT = {
        'reply': 'C-1001 케이스가 유사합니다.',
        'tool_calls': [{'name': 'search_cases', 'input': {'query': 'VRRP'}}],
        'model': 'claude-haiku-4-5-20251001',
        'agent': 'search',
    }

    def setUp(self):
        from django.contrib.auth.models import User
        from .permissions import set_user_role
        for username, role in (('v1', 'viewer'), ('e1', 'engineer'), ('e2', 'engineer')):
            user = User.objects.create_user(username, password='role-pass-123!')
            set_user_role(user, role)

    def login(self, username):
        self.client.post('/api/auth/login/',
                         {'username': username, 'password': 'role-pass-123!'},
                         content_type='application/json')

    def chat(self, content, session_id=None):
        with patch('api.views.help_agent.chat', return_value=dict(self.FAKE_RESULT)):
            return self.client.post(
                '/api/help-agent/chat/',
                {'messages': [{'role': 'user', 'content': content}],
                 'session_id': session_id},
                content_type='application/json')

    def test_engineer_can_chat_viewer_cannot(self):
        self.login('e1')
        self.assertEqual(self.chat('VRRP 유사 사례').status_code, 200)
        self.login('v1')
        self.assertEqual(self.chat('VRRP 유사 사례').status_code, 403)

    def test_first_chat_creates_session_with_both_turns(self):
        from .models import ChatSession, ChatTurn
        self.login('e1')
        data = self.chat('VRRP 유사 사례 찾아줘').json()

        session = ChatSession.objects.get(id=data['session_id'])
        self.assertEqual(session.user.username, 'e1')
        self.assertEqual(session.title, 'VRRP 유사 사례 찾아줘')
        turns = list(ChatTurn.objects.filter(session=session))
        self.assertEqual([t.role for t in turns], ['user', 'assistant'])
        self.assertEqual(turns[1].agent, 'search')
        self.assertEqual(turns[1].tool_calls, self.FAKE_RESULT['tool_calls'])

    def test_followup_appends_to_same_session(self):
        from .models import ChatTurn
        self.login('e1')
        first = self.chat('첫 질문').json()
        second = self.chat('두 번째 질문', session_id=first['session_id']).json()
        self.assertEqual(second['session_id'], first['session_id'])
        self.assertEqual(
            ChatTurn.objects.filter(session_id=first['session_id']).count(), 4)

    def test_sessions_are_private_to_owner(self):
        self.login('e1')
        session_id = self.chat('e1의 질문').json()['session_id']

        self.login('e2')
        self.assertEqual(
            self.client.get(f'/api/help-agent/sessions/{session_id}/').status_code, 404)
        self.assertEqual(self.client.get('/api/help-agent/sessions/').json(), [])
        # 남의 세션에 이어 쓰기도 차단
        self.assertEqual(self.chat('가로채기', session_id=session_id).status_code, 404)

    def test_session_list_and_detail_and_delete(self):
        from .models import ChatSession
        self.login('e1')
        session_id = self.chat('목록 테스트').json()['session_id']

        sessions = self.client.get('/api/help-agent/sessions/').json()
        self.assertEqual([s['id'] for s in sessions], [session_id])
        self.assertEqual(sessions[0]['turn_count'], 2)

        detail = self.client.get(f'/api/help-agent/sessions/{session_id}/').json()
        self.assertEqual(len(detail['turns']), 2)
        self.assertEqual(detail['turns'][1]['content'], self.FAKE_RESULT['reply'])

        delete = self.client.delete(f'/api/help-agent/sessions/{session_id}/')
        self.assertEqual(delete.status_code, 204)
        self.assertFalse(ChatSession.objects.filter(id=session_id).exists())

    def test_save_failure_still_returns_reply(self):
        self.login('e1')
        with patch('api.views.ChatSession.objects.create', side_effect=RuntimeError('db down')):
            response = self.chat('저장 실패해도 답변은 온다')
        self.assertEqual(response.status_code, 200)
        data = response.json()
        self.assertEqual(data['reply'], self.FAKE_RESULT['reply'])
        self.assertIsNone(data['session_id'])


class HelpAgentToolLoopTests(TestCase):
    """도구 호출 루프가 반복 상한을 소진해도 최종 답변을 생성해야 한다."""

    @staticmethod
    def _tool_response(i):
        return SimpleNamespace(
            stop_reason='tool_use',
            content=[SimpleNamespace(type='tool_use', name='search_cases',
                                     input={'query': f'q{i}'}, id=f'tool_{i}')],
        )

    def test_exhausted_tool_loop_forces_final_answer(self):
        final = SimpleNamespace(
            stop_reason='end_turn',
            content=[SimpleNamespace(type='text', text='수집한 근거 기준 최종 답변')],
        )
        calls = []

        def create(**kwargs):
            calls.append(kwargs)
            if kwargs.get('tool_choice') == {'type': 'none'}:
                return final
            return self._tool_response(len(calls))

        client = SimpleNamespace(messages=SimpleNamespace(create=create))
        reply, trace, _, _ = help_agent._run_agent(
            client, 'search', [{'role': 'user', 'content': '질문'}])

        self.assertEqual(reply, '수집한 근거 기준 최종 답변')
        # 루프 상한만큼 도구 호출 후, 마무리 호출은 도구를 차단한다
        self.assertEqual(len(trace), help_agent.MAX_TOOL_ITERATIONS)
        self.assertEqual(len(calls), help_agent.MAX_TOOL_ITERATIONS + 1)
        self.assertEqual(calls[-1]['tool_choice'], {'type': 'none'})

    def test_wrap_up_failure_falls_back_to_last_response(self):
        calls = []

        def create(**kwargs):
            calls.append(kwargs)
            if kwargs.get('tool_choice') == {'type': 'none'}:
                raise anthropic.APIConnectionError(request=MagicMock())
            return self._tool_response(len(calls))

        client = SimpleNamespace(messages=SimpleNamespace(create=create))
        reply, _, _, _ = help_agent._run_agent(
            client, 'search', [{'role': 'user', 'content': '질문'}])
        self.assertEqual(reply, '')  # chat()에서 안내 문구로 대체된다


class ChatAttachmentTests(TestCase):
    """첨부(멀티모달): 블록 변환·업로드 검증·저장·삭제."""

    IMAGE = {'file_id': 'file_img1', 'filename': 'err.png',
             'kind': 'image', 'size_bytes': 1024}
    PDF = {'file_id': 'file_doc1', 'filename': 'guide.pdf',
           'kind': 'document', 'size_bytes': 2048}

    def test_recent_attachment_becomes_file_block(self):
        messages = [{'role': 'user', 'content': '이 에러 뭔가요',
                     'attachments': [self.IMAGE]}]
        expanded, has_files = help_agent._expand_attachments(messages)

        self.assertTrue(has_files)
        blocks = expanded[0]['content']
        self.assertEqual(blocks[0], {'type': 'image',
                                     'source': {'type': 'file', 'file_id': 'file_img1'}})
        self.assertEqual(blocks[1]['type'], 'text')
        self.assertIn('err.png', blocks[1]['text'])

    def test_pdf_uses_document_block(self):
        expanded, _ = help_agent._expand_attachments(
            [{'role': 'user', 'content': '요약해줘', 'attachments': [self.PDF]}])
        self.assertEqual(expanded[0]['content'][0]['type'], 'document')

    def test_old_attachments_degrade_to_filename_note(self):
        """오래된 첨부는 파일을 다시 올리지 않는다 — 이력이 길어져도 재과금 없음."""
        messages = [
            {'role': 'user', 'content': f'질문{i}', 'attachments': [dict(
                self.IMAGE, file_id=f'file_{i}', filename=f'shot{i}.png')]}
            for i in range(help_agent.MAX_ATTACHMENT_MESSAGES + 2)
        ]
        expanded, _ = help_agent._expand_attachments(messages)

        live = [m for m in expanded if isinstance(m['content'], list)]
        self.assertEqual(len(live), help_agent.MAX_ATTACHMENT_MESSAGES)
        # 잘려나간 오래된 턴은 파일명 표시만 텍스트로 남는다
        self.assertIsInstance(expanded[0]['content'], str)
        self.assertIn('shot0.png', expanded[0]['content'])

    def test_messages_without_attachments_are_untouched(self):
        messages = [{'role': 'user', 'content': '평범한 질문'}]
        expanded, has_files = help_agent._expand_attachments(messages)
        self.assertFalse(has_files)
        self.assertEqual(expanded, messages)

    def test_question_text_falls_back_to_filenames(self):
        """본문 없이 스크린샷만 올려도 분류·제목이 비지 않아야 한다."""
        text = help_agent.question_text(
            {'role': 'user', 'content': '', 'attachments': [self.IMAGE]})
        self.assertIn('err.png', text)

    def test_upload_rejects_unsupported_extension(self):
        with self.assertRaises(help_agent.AttachmentRejected):
            help_agent.upload_attachment('config.exe', b'data')

    def test_upload_rejects_oversized_file(self):
        oversized = b'x' * (help_agent.MAX_ATTACHMENT_BYTES + 1)
        with self.assertRaises(help_agent.AttachmentRejected):
            help_agent.upload_attachment('shot.png', oversized)

    @override_settings(ANTHROPIC_API_KEY='test-key')
    def test_upload_returns_metadata(self):
        client = MagicMock()
        client.beta.files.upload.return_value = SimpleNamespace(id='file_new')
        with patch('api.services.help_agent.anthropic.Anthropic', return_value=client):
            meta = help_agent.upload_attachment('shot.PNG', b'binary')

        self.assertEqual(meta, {'file_id': 'file_new', 'filename': 'shot.PNG',
                                'kind': 'image', 'size_bytes': 6,
                                'converted': False})

    def test_attachment_beta_header_only_when_files_present(self):
        """첨부가 없으면 기존처럼 정식 엔드포인트를 쓴다 (베타 의존 최소화)."""
        beta_calls, plain_calls = [], []
        answer = SimpleNamespace(
            stop_reason='end_turn',
            content=[SimpleNamespace(type='text', text='답변')])
        client = SimpleNamespace(
            messages=SimpleNamespace(
                create=lambda **kw: (plain_calls.append(kw), answer)[1]),
            beta=SimpleNamespace(messages=SimpleNamespace(
                create=lambda **kw: (beta_calls.append(kw), answer)[1])))

        help_agent._run_agent(client, 'search', [{'role': 'user', 'content': '질문'}])
        self.assertEqual((len(plain_calls), len(beta_calls)), (1, 0))

        help_agent._run_agent(client, 'search', [
            {'role': 'user', 'content': '이 화면 봐줘', 'attachments': [self.IMAGE]}])
        self.assertEqual(len(beta_calls), 1)
        self.assertEqual(beta_calls[0]['betas'], [help_agent.FILES_BETA])


class ChatAttachmentApiTests(TestCase):
    """첨부가 붙은 채팅 요청의 검증·저장·정리."""

    FAKE_RESULT = {'reply': '답변', 'tool_calls': [], 'model': 'm', 'agent': 'search'}
    IMAGE = {'file_id': 'file_img1', 'filename': 'err.png',
             'kind': 'image', 'size_bytes': 1024}

    def setUp(self):
        from .permissions import set_user_role
        user = User.objects.create_user('e1', password='role-pass-123!')
        set_user_role(user, 'engineer')
        self.client.post('/api/auth/login/',
                         {'username': 'e1', 'password': 'role-pass-123!'},
                         content_type='application/json')

    def chat(self, message):
        with patch('api.views.help_agent.chat', return_value=dict(self.FAKE_RESULT)):
            return self.client.post('/api/help-agent/chat/',
                                    {'messages': [message]},
                                    content_type='application/json')

    def test_attachment_only_message_is_accepted_and_stored(self):
        from .models import ChatTurn
        response = self.chat({'role': 'user', 'content': '',
                              'attachments': [self.IMAGE]})
        self.assertEqual(response.status_code, 200)

        turn = ChatTurn.objects.get(session_id=response.json()['session_id'],
                                    role='user')
        self.assertEqual(turn.attachments, [self.IMAGE])

    def test_empty_message_without_attachment_is_rejected(self):
        response = self.chat({'role': 'user', 'content': '   '})
        self.assertEqual(response.status_code, 400)

    def test_malformed_attachment_is_rejected(self):
        response = self.chat({'role': 'user', 'content': '질문',
                              'attachments': [{'file_id': 'file_x', 'kind': 'video'}]})
        self.assertEqual(response.status_code, 400)

    def test_session_title_uses_filename_when_body_empty(self):
        from .models import ChatSession
        response = self.chat({'role': 'user', 'content': '',
                              'attachments': [self.IMAGE]})
        session = ChatSession.objects.get(id=response.json()['session_id'])
        self.assertIn('err.png', session.title)

    def test_deleting_session_removes_uploaded_files(self):
        session_id = self.chat({'role': 'user', 'content': '봐줘',
                                'attachments': [self.IMAGE]}).json()['session_id']

        with patch('api.views.help_agent.delete_files') as delete_files:
            response = self.client.delete(f'/api/help-agent/sessions/{session_id}/')

        self.assertEqual(response.status_code, 204)
        delete_files.assert_called_once_with(['file_img1'])

    def test_upload_endpoint_rejects_unsupported_type(self):
        upload = SimpleUploadedFile('payload.exe', b'binary')
        response = self.client.post('/api/help-agent/attachments/', {'file': upload})
        self.assertEqual(response.status_code, 400)

    def test_upload_endpoint_returns_file_id(self):
        upload = SimpleUploadedFile('shot.png', b'binary', content_type='image/png')
        with patch('api.views.help_agent.upload_attachment',
                   return_value=dict(self.IMAGE)) as uploader:
            response = self.client.post('/api/help-agent/attachments/', {'file': upload})

        self.assertEqual(response.status_code, 201)
        self.assertEqual(response.json()['file_id'], 'file_img1')
        self.assertEqual(uploader.call_args[0][0], 'shot.png')


class TriageRoutingTests(TestCase):
    """파일 생성 요청 라우팅 — 파일을 만들 수 있는 담당은 report뿐이다."""

    ATTACH = [{'file_id': 'file_x', 'filename': 'incidents.xlsx',
               'kind': 'document', 'size_bytes': 500}]

    @staticmethod
    def _client(verdict):
        """haiku 분류기가 verdict를 돌려주는 가짜 클라이언트."""
        response = SimpleNamespace(
            content=[SimpleNamespace(type='text', text=verdict)])
        return SimpleNamespace(
            messages=SimpleNamespace(create=lambda **kw: response))

    def triage(self, text, verdict='search', attachments=None):
        message = {'role': 'user', 'content': text}
        if attachments:
            message['attachments'] = attachments
        return help_agent._triage(self._client(verdict), [message])

    def test_file_output_requests_route_to_report_by_rule(self):
        """보고서 단어가 없어도 산출물 단어만으로 report여야 한다."""
        for text in ('이 엑셀 내용으로 PPT 만들어줘',
                     '첨부 시트를 슬라이드로 정리해줘',
                     '파워포인트로 바꿔줘',
                     '이거 발표자료로 만들어줘',
                     '첨부한 표를 워드 문서로 만들어줘'):
            # 분류기가 off_topic이라 답해도 규칙이 먼저 report로 잡아야 한다
            self.assertEqual(self.triage(text, verdict='off_topic'), 'report', text)

    def test_keyword_search_is_not_hijacked_by_word_substring(self):
        """'키워드'가 '워드'에 부분일치해 케이스 검색이 리포트로 새면 안 된다."""
        self.assertEqual(self.triage('이 키워드로 케이스 찾아줘'), 'search')
        self.assertEqual(self.triage('키워드 검색해줘'), 'search')

    def test_attachment_request_is_never_off_topic(self):
        """오분류로 고정 안내가 나가면 기능이 없는 것처럼 보인다."""
        self.assertEqual(
            self.triage('이거 좀 봐줘', verdict='off_topic', attachments=self.ATTACH),
            'search')

    def test_off_topic_still_works_without_attachment(self):
        """첨부가 없으면 기존대로 비용 가드가 동작해야 한다."""
        self.assertEqual(self.triage('오늘 점심 뭐 먹지', verdict='off_topic'),
                         'off_topic')

    def test_reading_an_attachment_stays_on_search(self):
        """파일 생성 요청이 아니면 굳이 리포팅으로 보내지 않는다."""
        self.assertEqual(
            self.triage('엑셀 내용 요약해줘', verdict='search', attachments=self.ATTACH),
            'search')


class FinalTextTests(TestCase):
    """코드 실행 중간 메모가 답변 앞에 붙지 않아야 한다 (실측: 영어 작업 로그 노출)."""

    @staticmethod
    def blocks(*specs):
        return [SimpleNamespace(type=t, text=v) for t, v in specs]

    def test_narration_before_tool_result_is_dropped(self):
        content = self.blocks(
            ('text', 'Content confirmed correct. Now visual QA.'),
            ('bash_code_execution_tool_result', None),
            ('text', '보고자료를 생성했습니다.'))
        self.assertEqual(help_agent._final_text(content), '보고자료를 생성했습니다.')

    def test_plain_answer_is_unchanged(self):
        content = self.blocks(('text', 'VRRP 유사 사례는 C-1122입니다.'))
        self.assertEqual(help_agent._final_text(content),
                         'VRRP 유사 사례는 C-1122입니다.')

    def test_consecutive_trailing_blocks_are_joined(self):
        content = self.blocks(
            ('bash_code_execution_tool_result', None),
            ('text', '앞부분 '), ('text', '뒷부분'))
        self.assertEqual(help_agent._final_text(content), '앞부분 뒷부분')

    def test_text_only_before_tool_result_is_kept_as_fallback(self):
        """도구 결과로 끝나면 앞의 텍스트라도 답변으로 살린다 (빈 답변 방지)."""
        content = self.blocks(('text', '설명만 있고 끝'),
                              ('bash_code_execution_tool_result', None))
        self.assertEqual(help_agent._final_text(content), '설명만 있고 끝')


class ReportPromptCostTests(TestCase):
    """리포팅 프롬프트가 비용·출처 규칙을 담고 있는지 (실측 $9.03 → $0.77 근거)."""

    def test_image_rendering_is_forbidden(self):
        """'한 번만'식 완화 지시로는 비용이 안 잡혔다 ($2.34~$12.18 편차) —
        이미지 변환 자체를 금지해야 한다."""
        prompt = help_agent.REPORT_SYSTEM_PROMPT
        self.assertIn('이미지로 변환해 눈으로 확인하지 마세요', prompt)
        self.assertIn('코드로 값을 재세요', prompt)

    def test_attachment_is_declared_as_the_data_source(self):
        """첨부로 문서를 만들 때 케이스 DB 수치가 섞이면 틀린 문서가 나간다."""
        prompt = help_agent.REPORT_SYSTEM_PROMPT
        self.assertIn('첨부', prompt)
        self.assertNotIn('반드시 get_case_stats', prompt)


class OfficeAttachmentTests(TestCase):
    """워드·엑셀·PPT 첨부: 서버에서 본문 텍스트로 변환해 올린다.

    Anthropic document 블록이 오피스 형식을 거부하고(400), 코드 실행 컨테이너로
    읽히기는 하나 파일당 ~2만 토큰이 들어서 택한 방식.
    """

    @staticmethod
    def make_docx(paragraphs=(), table_rows=()):
        from docx import Document
        document = Document()
        for text in paragraphs:
            document.add_paragraph(text)
        if table_rows:
            table = document.add_table(rows=0, cols=len(table_rows[0]))
            for values in table_rows:
                cells = table.add_row().cells
                for cell, value in zip(cells, values):
                    cell.text = value
        buffer = BytesIO()
        document.save(buffer)
        return buffer.getvalue()

    @staticmethod
    def make_xlsx(sheets):
        from openpyxl import Workbook
        workbook = Workbook()
        workbook.remove(workbook.active)
        for title, rows in sheets.items():
            sheet = workbook.create_sheet(title)
            for row in rows:
                sheet.append(row)
        buffer = BytesIO()
        workbook.save(buffer)
        return buffer.getvalue()

    @staticmethod
    def make_pptx(slides):
        """slides: [{'texts': [...], 'table': [[...]], 'notes': '...', 'group': [...]}]"""
        from pptx import Presentation
        from pptx.util import Inches
        presentation = Presentation()
        for spec in slides:
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            for text in spec.get('texts', ()):
                box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
                box.text_frame.text = text
            if spec.get('group'):
                group = slide.shapes.add_group_shape()
                for text in spec['group']:
                    box = group.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
                    box.text_frame.text = text
            rows = spec.get('table')
            if rows:
                table = slide.shapes.add_table(len(rows), len(rows[0]), Inches(1),
                                               Inches(3), Inches(4), Inches(1)).table
                for r, values in enumerate(rows):
                    for c, value in enumerate(values):
                        table.cell(r, c).text = value
            if spec.get('notes'):
                slide.notes_slide.notes_text_frame.text = spec['notes']
        buffer = BytesIO()
        presentation.save(buffer)
        return buffer.getvalue()

    def test_pptx_labels_each_slide_in_order(self):
        """슬라이드 경계가 없으면 도형에 흩어진 본문이 뒤섞여 맥락을 잃는다."""
        data = self.make_pptx([{'texts': ['장애 개요']}, {'texts': ['조치 결과']}])
        text = help_agent._extract_pptx(data)

        self.assertIn('[슬라이드 1]', text)
        self.assertIn('[슬라이드 2]', text)
        self.assertLess(text.index('장애 개요'), text.index('조치 결과'))

    def test_pptx_reads_tables_groups_and_speaker_notes(self):
        """그룹 도형은 재귀로 들어가야 안쪽 텍스트가 빠지지 않는다."""
        data = self.make_pptx([{
            'group': ['그룹 안 원인 설명'],
            'table': [['Bug', 'Fixed'], ['ACOS-104904', '6.0.9']],
            'notes': '고객 보고용 배경',
        }])
        text = help_agent._extract_pptx(data)

        self.assertIn('그룹 안 원인 설명', text)
        self.assertIn('ACOS-104904 | 6.0.9', text)
        self.assertIn('(발표자 노트) 고객 보고용 배경', text)

    def test_image_only_pptx_is_rejected_with_guidance(self):
        """도형만 있고 텍스트가 없는 발표자료는 캡처해 올리라고 안내한다."""
        with self.assertRaises(help_agent.AttachmentRejected) as caught:
            help_agent._office_to_text('.pptx', 'deck.pptx', self.make_pptx([{}]))
        self.assertIn('캡처', str(caught.exception))

    @override_settings(ANTHROPIC_API_KEY='test-key')
    def test_pptx_upload_is_converted_to_text(self):
        original = self.make_pptx([{'texts': ['VRRP 이중화 점검 결과']}])
        client = MagicMock()
        client.beta.files.upload.return_value = SimpleNamespace(id='file_ppt')
        with patch('api.services.help_agent.anthropic.Anthropic', return_value=client):
            meta = help_agent.upload_attachment('deck.pptx', original)

        self.assertTrue(meta['converted'])
        self.assertEqual(meta['filename'], 'deck.pptx')  # 사용자에겐 원본 이름
        name, stream, mime = client.beta.files.upload.call_args.kwargs['file']
        self.assertEqual((name, mime), ('deck.pptx.txt', 'text/plain'))
        self.assertIn('VRRP 이중화 점검 결과', stream.getvalue().decode())

    @override_settings(ANTHROPIC_API_KEY='test-key')
    def test_legacy_office_upload_names_the_format_to_save_as(self):
        """"PPT 첨부"를 시도하면 .ppt를 올리기 쉽다 — 형식 나열 대신 방법을 알려준다."""
        for name, modern in (('deck.ppt', '.pptx'), ('report.doc', '.docx'),
                             ('sheet.xls', '.xlsx')):
            with self.assertRaises(help_agent.AttachmentRejected) as caught:
                help_agent.upload_attachment(name, b'ole2-binary')
            self.assertIn(modern, str(caught.exception))

    def test_docx_keeps_document_order_of_paragraphs_and_tables(self):
        """표를 뒤로 몰아버리면 맥락이 끊긴다 — 원래 순서를 지켜야 한다."""
        data = self.make_docx(paragraphs=['머리말', '꼬리말'],
                              table_rows=[['Bug', 'Fixed'], ['ACOS-104904', '6.0.9']])
        text = help_agent._extract_docx(data)

        self.assertIn('머리말', text)
        self.assertIn('ACOS-104904 | 6.0.9', text)
        # add_table은 문단 뒤에 붙으므로 꼬리말보다 표가 뒤에 와야 정상
        self.assertLess(text.index('머리말'), text.index('ACOS-104904'))

    def test_xlsx_labels_each_sheet_and_skips_empty_rows(self):
        data = self.make_xlsx({
            'Issues': [['Bug ID', 'Fixed In'], [], ['ACOS-104904', '6.0.9']],
            'Notes': [['비고', '없음']],
        })
        text = help_agent._extract_xlsx(data)

        self.assertIn('[시트: Issues]', text)
        self.assertIn('[시트: Notes]', text)
        self.assertIn('ACOS-104904 | 6.0.9', text)
        self.assertNotIn('\n\n', text)  # 빈 행이 그대로 실리지 않는다

    def test_conversion_adds_source_header(self):
        """모델이 원본이 무엇이었는지 알아야 한다."""
        data = help_agent._office_to_text(
            '.docx', 'report.docx', self.make_docx(paragraphs=['본문']))
        self.assertTrue(data.decode().startswith('[report.docx에서 추출한 텍스트]'))

    def test_text_only_document_is_rejected_with_guidance(self):
        """이미지로만 된 문서는 캡처해서 올리라고 안내해야 한다."""
        with self.assertRaises(help_agent.AttachmentRejected) as caught:
            help_agent._office_to_text('.docx', 'empty.docx', self.make_docx())
        self.assertIn('캡처', str(caught.exception))

    def test_corrupt_file_is_rejected_not_crashed(self):
        with self.assertRaises(help_agent.AttachmentRejected) as caught:
            help_agent._office_to_text('.xlsx', 'broken.xlsx', b'not a zip')
        self.assertIn('읽을 수 없습니다', str(caught.exception))

    def test_oversized_extraction_is_truncated_with_notice(self):
        """행이 수천 개인 시트를 통째로 실으면 질문 한 번에 수만 토큰이 나간다."""
        rows = [[f'row-{i}' * 20] for i in range(3000)]
        data = help_agent._office_to_text(
            '.xlsx', 'big.xlsx', self.make_xlsx({'Sheet': rows}))
        text = data.decode()

        self.assertIn('이하 생략', text)
        self.assertLess(len(text), help_agent.MAX_EXTRACTED_CHARS + 500)

    @override_settings(ANTHROPIC_API_KEY='test-key')
    def test_upload_sends_text_but_reports_original_file(self):
        """사용자에겐 원래 파일명·크기를, 모델에겐 추출 텍스트를 준다."""
        original = self.make_docx(paragraphs=['VRRP failover 원인 분석'])
        client = MagicMock()
        client.beta.files.upload.return_value = SimpleNamespace(id='file_conv')
        with patch('api.services.help_agent.anthropic.Anthropic', return_value=client):
            meta = help_agent.upload_attachment('report.docx', original)

        self.assertEqual(meta, {'file_id': 'file_conv', 'filename': 'report.docx',
                                'kind': 'document', 'size_bytes': len(original),
                                'converted': True})
        name, stream, mime = client.beta.files.upload.call_args.kwargs['file']
        self.assertEqual((name, mime), ('report.docx.txt', 'text/plain'))
        self.assertIn('VRRP failover 원인 분석', stream.getvalue().decode())

    @override_settings(ANTHROPIC_API_KEY='test-key')
    def test_non_office_upload_is_not_converted(self):
        client = MagicMock()
        client.beta.files.upload.return_value = SimpleNamespace(id='file_raw')
        with patch('api.services.help_agent.anthropic.Anthropic', return_value=client):
            meta = help_agent.upload_attachment('shot.png', b'binary')

        self.assertFalse(meta['converted'])
        name, stream, mime = client.beta.files.upload.call_args.kwargs['file']
        self.assertEqual((name, mime, stream.getvalue()), ('shot.png', 'image/png', b'binary'))


class PurgeChatAttachmentsTests(TestCase):
    """고아 첨부 정리: 참조된 파일(템플릿·생성 문서·대화 첨부)은 절대 지우지 않는다."""

    def setUp(self):
        from .models import AppSetting, ChatSession, ChatTurn
        session = ChatSession.objects.create(title='대화')
        ChatTurn.objects.create(
            session=session, role='user', content='봐줘',
            attachments=[{'file_id': 'file_attached', 'filename': 'a.png',
                          'kind': 'image', 'size_bytes': 1}])
        ChatTurn.objects.create(
            session=session, role='assistant', content='답변',
            files=[{'file_id': 'file_report', 'filename': 'r.docx', 'size_bytes': 2}])
        AppSetting.set('report_template_docx', 'abc123:file_template')

    @staticmethod
    def _file(file_id, days_old):
        return SimpleNamespace(id=file_id, filename=f'{file_id}.png',
                               created_at=timezone.now() - timedelta(days=days_old))

    def run_command(self, files, **options):
        client = MagicMock()
        client.beta.files.list.return_value = files
        out = StringIO()
        with override_settings(ANTHROPIC_API_KEY='test-key'), \
                patch('anthropic.Anthropic', return_value=client):
            call_command('purge_chat_attachments', stdout=out, **options)
        return client, out.getvalue()

    def test_referenced_files_are_never_listed(self):
        client, output = self.run_command([
            self._file('file_attached', 30),
            self._file('file_report', 30),
            self._file('file_template', 30),
            self._file('file_orphan', 30),
        ], apply=True)

        client.beta.files.delete.assert_called_once_with('file_orphan')
        self.assertNotIn('file_template', output)

    def test_recent_orphans_are_kept_during_grace_period(self):
        """업로드만 하고 아직 질문을 안 보낸 파일을 지워버리면 안 된다."""
        client, _ = self.run_command([self._file('file_fresh', 0)], apply=True)
        client.beta.files.delete.assert_not_called()

    def test_dry_run_is_the_default(self):
        client, output = self.run_command([self._file('file_orphan', 30)])
        client.beta.files.delete.assert_not_called()
        self.assertIn('file_orphan', output)
        self.assertIn('--apply', output)


class FetchUrlToolTests(TestCase):
    """fetch_url 도구: SSRF 차단과 본문 추출."""

    def test_non_http_scheme_blocked(self):
        out = json.loads(help_agent._fetch_url('ftp://example.com/file'))
        self.assertIn('http/https', out['error'])

    def test_private_ip_blocked(self):
        out = json.loads(help_agent._fetch_url('http://192.168.74.158/admin'))
        self.assertIn('내부망', out['error'])

    def test_localhost_blocked(self):
        out = json.loads(help_agent._fetch_url('http://localhost:8000/api/'))
        self.assertIn('내부망', out['error'])

    def test_html_content_extracted(self):
        html = (b'<html><head><title>Security Advisory</title></head>'
                b'<body><script>tracker()</script>'
                b'<p>ACOS 6.0.9 fixes CVE-2026-45447</p></body></html>')
        title, text = help_agent._extract_page_content('text/html; charset=utf-8', html)
        self.assertEqual(title, 'Security Advisory')
        self.assertIn('ACOS 6.0.9', text)
        self.assertNotIn('tracker()', text)

    def test_plain_text_passthrough(self):
        title, text = help_agent._extract_page_content('text/plain', 'release note 본문'.encode())
        self.assertEqual(text, 'release note 본문')

    def test_fetch_returns_payload_and_truncates(self):
        body = (b'<html><head><title>Long</title></head><body>'
                + b'A' * (help_agent.FETCH_MAX_CHARS + 100) + b'</body></html>')
        response = MagicMock()
        response.headers = {'content-type': 'text/html'}
        response.iter_bytes.return_value = [body]
        response.url = 'https://vendor.example/advisory'
        client = MagicMock()
        client.stream.return_value.__enter__.return_value = response

        with patch.object(help_agent.httpx, 'Client') as client_cls, \
             patch.object(help_agent, '_assert_public_http_url'):
            client_cls.return_value.__enter__.return_value = client
            out = json.loads(help_agent._fetch_url('https://vendor.example/advisory'))

        self.assertEqual(out['title'], 'Long')
        self.assertEqual(len(out['content']), help_agent.FETCH_MAX_CHARS)
        self.assertIn('notice', out)


class ChatKnowledgeExtractTests(TestCase):
    """대화 세션 -> 지식 추출 (2단계): 명시적 버튼, AI 정제, draft 등록."""

    EXTRACTED = {
        'has_knowledge': True,
        'vendor': 'A10',
        'title': 'VRRP failover 시 세션 동기화 누락',
        'problem': 'failover 후 기존 세션이 끊깁니다.',
        'root_cause': 'session sync 미설정.',
        'resolution': 'vrrp-a session-sync enable 설정을 추가합니다.',
        'device_model': 'TH4435',
        'software_version': '5.2.1-P10',
    }

    def setUp(self):
        from django.contrib.auth.models import User
        from .permissions import set_user_role
        from .models import ChatSession, ChatTurn
        for username in ('e1', 'e2'):
            user = User.objects.create_user(username, password='role-pass-123!')
            set_user_role(user, 'engineer')
        owner = User.objects.get(username='e1')
        self.session = ChatSession.objects.create(user=owner, title='VRRP 문제')
        ChatTurn.objects.create(session=self.session, role='user',
                                content='VRRP failover 후 세션이 끊겨요')
        ChatTurn.objects.create(
            session=self.session, role='assistant', agent='tech',
            content='vrrp-a session-sync enable 설정이 필요합니다.',
            tool_calls=[{'name': 'search_references', 'input': {'query': 'vrrp'}}])

    def login(self, username):
        self.client.post('/api/auth/login/',
                         {'username': username, 'password': 'role-pass-123!'},
                         content_type='application/json')

    def extract(self, session_id=None, ai_result='default'):
        if ai_result == 'default':
            ai_result = dict(self.EXTRACTED)
        with patch('api.services.knowledge.generate_structured_with_model',
                   return_value=(MODEL, ai_result)), \
             patch('api.services.knowledge.enrich_with_references',
                   return_value='no_candidates'):
            return self.client.post(
                f'/api/help-agent/sessions/{session_id or self.session.id}/knowledge/')

    def test_extracts_draft_knowledge_with_session_source(self):
        from .models import KnowledgeItem
        self.login('e1')
        res = self.extract()
        self.assertEqual(res.status_code, 201)
        data = res.json()
        self.assertEqual(data['outcome'], 'created')
        self.assertEqual(data['item']['source_session']['id'], self.session.id)
        self.assertIsNone(data['item']['source_case'])

        item = KnowledgeItem.objects.get(id=data['item']['id'])
        self.assertEqual(item.chat_session, self.session)
        self.assertEqual(item.vendor, 'A10')
        self.assertEqual(item.status, 'draft')
        self.assertIn('session-sync', item.resolution)

    def test_second_extract_returns_existing(self):
        from .models import KnowledgeItem
        self.login('e1')
        first = self.extract().json()
        res = self.extract()
        self.assertEqual(res.status_code, 200)
        self.assertEqual(res.json()['outcome'], 'exists')
        self.assertEqual(res.json()['item']['id'], first['item']['id'])
        self.assertEqual(KnowledgeItem.objects.count(), 1)

    def test_no_knowledge_conversation_rejected(self):
        self.login('e1')
        res = self.extract(ai_result={**self.EXTRACTED, 'has_knowledge': False,
                                      'resolution': ''})
        self.assertEqual(res.status_code, 400)
        self.assertEqual(res.json()['outcome'], 'no_knowledge')

    def test_unknown_vendor_rejected(self):
        self.login('e1')
        res = self.extract(ai_result={**self.EXTRACTED, 'vendor': 'Unknown'})
        self.assertEqual(res.status_code, 400)
        self.assertEqual(res.json()['outcome'], 'no_vendor')

    def test_chat_schema_has_no_empty_enum_value(self):
        # Gemini는 enum의 빈 문자열을 400 INVALID_ARGUMENT로 거부한다
        from .services.knowledge import CHAT_KNOWLEDGE_SCHEMA
        self.assertNotIn('', CHAT_KNOWLEDGE_SCHEMA['properties']['vendor']['enum'])

    def test_other_users_session_not_found(self):
        self.login('e2')
        self.assertEqual(self.extract().status_code, 404)

    def test_chat_material_includes_tools_and_roles(self):
        from .services.knowledge import build_chat_material
        material = build_chat_material(self.session)
        self.assertIn('[엔지니어]', material)
        self.assertIn('[AI] (사용 도구: search_references)', material)
        self.assertIn('vrrp-a session-sync enable', material)


class KnowledgeSyncTests(TestCase):
    """지식 동기화 버튼: 미검토 Resolved 케이스 일괄 추출 (admin 전용)."""

    EXTRACTED = {
        'has_knowledge': True,
        'title': '문제 요약',
        'problem': '증상',
        'root_cause': '원인',
        'resolution': 'fix 명령을 실행합니다.',
        'device_model': '',
        'software_version': '',
    }

    def setUp(self):
        from django.contrib.auth.models import User
        from .permissions import set_user_role
        for username, role in (('e1', 'engineer'), ('a1', 'admin')):
            user = User.objects.create_user(username, password='role-pass-123!')
            set_user_role(user, role)

    def login(self, username):
        self.client.post('/api/auth/login/',
                         {'username': username, 'password': 'role-pass-123!'},
                         content_type='application/json')

    def sync(self, ai_result='default'):
        if ai_result == 'default':
            ai_result = dict(self.EXTRACTED)
        with patch('api.services.knowledge.generate_structured_with_model',
                   return_value=(MODEL, ai_result)), \
             patch('api.services.knowledge.enrich_with_references',
                   return_value='no_candidates'):
            return self.client.post('/api/knowledge/sync/')

    def test_engineer_is_blocked(self):
        self.login('e1')
        self.assertEqual(self.sync().status_code, 403)

    def test_sync_extracts_and_marks_checked(self):
        from .models import KnowledgeItem
        resolved = make_case(vendor='A10', status='Resolved', summary='해결된 케이스')
        make_case(vendor='A10', status='Open', summary='미해결 케이스')

        self.login('a1')
        data = self.sync().json()
        self.assertEqual(data, {'scanned': 1, 'created': 1, 'no_knowledge': 0,
                                'failed': 0, 'remaining': 0})
        item = KnowledgeItem.objects.get()
        self.assertEqual(item.case, resolved)
        resolved.refresh_from_db()
        self.assertIsNotNone(resolved.knowledge_checked_at)

    def test_no_knowledge_case_is_not_rescanned(self):
        case = make_case(vendor='A10', status='Resolved', summary='공지 케이스')
        self.login('a1')
        data = self.sync(ai_result={**self.EXTRACTED, 'has_knowledge': False,
                                    'resolution': ''}).json()
        self.assertEqual((data['scanned'], data['no_knowledge']), (1, 1))
        case.refresh_from_db()
        self.assertIsNotNone(case.knowledge_checked_at)
        # 재클릭: 이미 검토된 케이스는 다시 스캔하지 않는다 (AI 비용 절감)
        data = self.sync().json()
        self.assertEqual(data['scanned'], 0)

    def test_failed_case_remains_for_retry(self):
        make_case(vendor='A10', status='Resolved', summary='AI 오류 케이스')
        self.login('a1')
        data = self.sync(ai_result=None).json()
        self.assertEqual((data['failed'], data['remaining']), (1, 1))
        # 실패 건은 검토 표시가 없어 다음 동기화에서 재시도된다
        data = self.sync().json()
        self.assertEqual((data['scanned'], data['created']), (1, 1))


class SignupRequestTests(TestCase):
    """계정 발급 — 신청 즉시 생성 + 관리자 알림 (2026-08-11 승인 방식 폐기).

    승인 링크(GET)는 메일 보안 스캐너가 사람보다 먼저 눌러 자동 승인시키는
    문제가 있어 제거했다. 통제는 사전 승인 대신 사후 조치(계정 관리)로 옮겼다.
    """

    def request_signup(self, **overrides):
        data = {'username': 'newbie', 'password': 'newbie-pass-77!',
                'name': '신입', 'email': 'newbie@ubersys.co.kr',
                'reason': '케이스 조회 필요'}
        data.update(overrides)
        with patch('api.auth_views.gmail_client.send_email') as mock_send:
            response = self.client.post('/api/auth/signup-requests/', data,
                                        content_type='application/json')
        return response, mock_send

    def test_account_is_usable_immediately(self):
        from django.contrib.auth.models import User
        from api.permissions import get_user_role
        response, _ = self.request_signup()
        self.assertEqual(response.status_code, 201)

        user = User.objects.get(username='newbie')
        self.assertEqual(get_user_role(user), 'viewer')
        self.assertEqual(user.email, 'newbie@ubersys.co.kr')
        # 승인 절차 없이 바로 로그인된다
        login = self.client.post('/api/auth/login/',
                                 {'username': 'newbie', 'password': 'newbie-pass-77!'},
                                 content_type='application/json')
        self.assertEqual(login.status_code, 200)

    def test_admin_is_notified_without_password(self):
        response, mock_send = self.request_signup()
        self.assertEqual(response.status_code, 201)
        mock_send.assert_called_once()
        to, subject, html = mock_send.call_args[0]
        self.assertEqual(to, 'jhshin@ubersys.co.kr')
        self.assertIn('새 사용자 가입', subject)
        self.assertIn('newbie', html)
        self.assertIn('newbie@ubersys.co.kr', html)
        self.assertNotIn('newbie-pass-77!', html)  # 비밀번호는 메일에 없음
        # 승인 링크가 더 이상 없어야 한다 (스캐너 자동 클릭 원인 제거)
        self.assertNotIn('signup-approve', html)

    def test_requested_role_engineer_is_granted(self):
        from django.contrib.auth.models import User
        from api.permissions import get_user_role
        self.request_signup(requested_role='engineer')
        self.assertEqual(get_user_role(User.objects.get(username='newbie')), 'engineer')

    def test_requesting_admin_role_is_rejected(self):
        from django.contrib.auth.models import User
        response, mock_send = self.request_signup(requested_role='admin')
        self.assertEqual(response.status_code, 400)
        mock_send.assert_not_called()
        self.assertFalse(User.objects.filter(username='newbie').exists())

    def test_duplicate_username_rejected(self):
        from django.contrib.auth.models import User
        User.objects.create_user('taken', password='x-pass-123!')
        response, _ = self.request_signup(username='taken')
        self.assertEqual(response.status_code, 400)

        self.request_signup()
        response, _ = self.request_signup()  # 같은 아이디 재가입
        self.assertEqual(response.status_code, 400)
        self.assertEqual(User.objects.filter(username='newbie').count(), 1)

    def test_email_is_required_and_validated(self):
        from django.contrib.auth.models import User
        for bad in ('', '골뱅이없음', 'a@b'):
            response, _ = self.request_signup(email=bad)
            self.assertEqual(response.status_code, 400, msg=bad)
        self.assertFalse(User.objects.filter(username='newbie').exists())

    def test_notification_failure_does_not_block_signup(self):
        """알림은 부가 기능 — 메일이 실패해도 계정은 살아 있어야 한다."""
        from django.contrib.auth.models import User
        with patch('api.auth_views.gmail_client.send_email',
                   side_effect=Exception('smtp down')):
            response = self.client.post(
                '/api/auth/signup-requests/',
                {'username': 'newbie', 'password': 'newbie-pass-77!',
                 'email': 'newbie@ubersys.co.kr'},
                content_type='application/json')
        self.assertEqual(response.status_code, 201)
        self.assertTrue(User.objects.filter(username='newbie').exists())


class CaseRelationTests(TestCase):
    """케이스 간 상호 참조 추가/해제."""

    def setUp(self):
        from django.contrib.auth.models import User
        from .permissions import set_user_role
        for username, role in (('rv1', 'viewer'), ('re1', 'engineer')):
            user = User.objects.create_user(username, password='rel-pass-123!')
            set_user_role(user, role)
        self.a = make_case(vendor='A10', summary='본 케이스입니다 다섯자이상')
        self.b = make_case(vendor='A10', summary='관련 케이스입니다 다섯자이상')

    def login(self, username):
        self.client.post('/api/auth/login/',
                         {'username': username, 'password': 'rel-pass-123!'},
                         content_type='application/json')

    def test_add_relation_by_display_number_is_symmetric(self):
        self.login('re1')
        response = self.client.post(f'/api/cases/{self.a.id}/relations/',
                                    {'case_id': f'C-{1000 + self.b.id}'},
                                    content_type='application/json')
        self.assertEqual(response.status_code, 201)
        # 양방향 반영 + 상세 응답에 포함
        detail_b = self.client.get(f'/api/cases/{self.b.id}/').json()
        self.assertEqual(detail_b['related_cases'][0]['id'], self.a.id)

    def test_remove_relation(self):
        self.login('re1')
        self.a.related_cases.add(self.b)
        response = self.client.delete(f'/api/cases/{self.a.id}/relations/{self.b.id}/')
        self.assertEqual(response.status_code, 200)
        self.assertEqual(self.a.related_cases.count(), 0)

    def test_viewer_cannot_modify_relations(self):
        self.login('rv1')
        response = self.client.post(f'/api/cases/{self.a.id}/relations/',
                                    {'case_id': f'C-{1000 + self.b.id}'},
                                    content_type='application/json')
        self.assertEqual(response.status_code, 403)

    def test_invalid_ref_and_self_ref_rejected(self):
        self.login('re1')
        bad = self.client.post(f'/api/cases/{self.a.id}/relations/',
                               {'case_id': 'C-9999'}, content_type='application/json')
        self.assertEqual(bad.status_code, 400)
        own = self.client.post(f'/api/cases/{self.a.id}/relations/',
                               {'case_id': f'C-{1000 + self.a.id}'},
                               content_type='application/json')
        self.assertEqual(own.status_code, 400)


@override_settings(GROUP_VENDOR_HINTS={'adc@ubersys.co.kr': 'A10'},
                   GMAIL_SYNC_INCLUDE_SUBJECTS=['Caseopen'])
class CustomerThreadVendorTests(TestCase):
    """벤더 도메인이 없는 고객사↔당사 스레드([Caseopen])의 벤더 추정."""

    def test_customer_mail_with_group_cc_gets_hinted_vendor(self):
        vendor, direction = detect_vendor_and_direction(
            '"엄현식" <hyunsik.um@samsung.com>',
            '"성의제" <ujseong22@ubersys.co.kr>',
            cc='"위버시스템즈(A10)" <adc@ubersys.co.kr>, IaaS NW <iaas.nw@samsung.com>',
        )
        self.assertEqual((vendor, direction), ('A10', 'inbound'))

    def test_our_reply_to_customer_is_outbound(self):
        vendor, direction = detect_vendor_and_direction(
            '"성의제" <ujseong22@ubersys.co.kr>',
            'hyunsik.um@samsung.com',
            cc='"adc@ubersys.co.kr" <adc@ubersys.co.kr>',
        )
        self.assertEqual((vendor, direction), ('A10', 'outbound'))

    def test_vendor_domain_still_wins_over_group_hint(self):
        vendor, direction = detect_vendor_and_direction(
            'tac@arista.com',
            'eng@ubersys.co.kr',
            cc='adc@ubersys.co.kr',
        )
        self.assertEqual((vendor, direction), ('Arista', 'inbound'))

    def test_no_hint_and_no_vendor_domain_returns_none(self):
        vendor, direction = detect_vendor_and_direction(
            'someone@samsung.com', 'eng@ubersys.co.kr', cc='other@ubersys.co.kr')
        self.assertIsNone(vendor)

    def test_gmail_query_includes_subject_keywords(self):
        query = build_gmail_query()
        self.assertIn('subject:Caseopen', query)
        # OR 그룹({}) 안에 들어가야 벤더 도메인 조건과 합집합이 된다
        self.assertIn('subject:Caseopen', query.split('}')[0])


class IgnoreRuleTests(TestCase):
    """공지/자동발송 메일이 케이스로 등록되는 것을 막는 규칙.

    실제로 쓰레기 케이스로 등록됐던 메일(Arista 공지 피드, Arista Community
    Central, HPE 계정 안내)의 발신자·제목을 그대로 사용한다.
    """

    GROUP_SENDER = '"\'Arista Networks\' via 기술부" <support@ubersys.co.kr>'

    def test_arista_notification_feed_subjects_are_ignored(self):
        for subject in ('New End of Sale email notification',
                        'New Field notice email notification',
                        'Security advisory Update email notification',
                        'Field notice Update email notification',
                        'New Software Release email notification'):
            reason = find_ignore_reason(self.GROUP_SENDER, subject)
            self.assertIsNotNone(reason, subject)

    def test_no_reply_sender_is_ignored(self):
        self.assertIsNotNone(find_ignore_reason(
            'Arista Community Central <no-reply@arista.com>',
            'Action Required: Please Update your Arista Community Central Nickname'))
        self.assertIsNotNone(find_ignore_reason(
            '"Hewlett Packard Enterprise (HPE)" <no-reply@auth.hpe.com>',
            'Action Required: Password Reset'))

    def test_relayed_no_reply_original_sender_is_ignored(self):
        # 그룹 중계로 From이 그룹 주소가 되어도 X-Original-Sender로 걸러진다
        reason = find_ignore_reason(
            self.GROUP_SENDER, 'Some vendor announcement',
            original_sender='noreply@arista.com')
        self.assertIsNotNone(reason)

    def test_real_case_mail_is_not_ignored(self):
        self.assertIsNone(find_ignore_reason(
            'A10 Customer Support Team <support@a10networks.com>',
            'A10 Networks Case Confirmation: NHN Cloud opened Case # 00457396'))
        self.assertIsNone(find_ignore_reason(
            self.GROUP_SENDER,
            'Re: New UBER Systems Co. Ltd Case: SR 834065 40G Interface Link FLAP',
            original_sender='tac-engineer@arista.com'))

    def test_gmail_query_excludes_notification_subjects(self):
        query = build_gmail_query()
        self.assertIn('-{', query)
        self.assertIn('subject:"email notification"', query.split('-{', 1)[1])


@override_settings(GMAIL_SYNC_INCLUDE_SUBJECTS=['Caseopen'])
class ExactSubjectMatchTests(TestCase):
    """스레드를 끊는 메일러(삼성 RE:(2) 카운터)의 케이스 오픈 스레드 병합."""

    SUBJECT = '[Caseopen] 수원 SCPv2 Multi-AZ 개발계 DATALB 파티션 변경 오류(API)'

    def test_clean_subject_strips_reply_counters(self):
        self.assertEqual(clean_subject(f'RE:(2) (2) {self.SUBJECT}'), self.SUBJECT)
        self.assertEqual(clean_subject(f'Re: (2) {self.SUBJECT}'), self.SUBJECT)

    def test_broken_thread_reply_matches_original_case(self):
        case = make_case(vendor='A10')
        make_email(case, self.SUBJECT, thread_id='thread-1')
        found = _find_case(None, 'thread-2', 'A10', f'RE:(2) (2) {self.SUBJECT}')
        self.assertEqual(found, case)

    def test_subject_without_open_keyword_is_not_merged(self):
        case = make_case(vendor='A10')
        make_email(case, '수원 SCPv2 개발계 정기 점검 안내', thread_id='thread-1')
        found = _find_case(None, 'thread-2', 'A10', '수원 SCPv2 개발계 정기 점검 안내')
        self.assertIsNone(found)

    def test_other_vendor_same_subject_is_not_merged(self):
        case = make_case(vendor='Arista')
        make_email(case, self.SUBJECT, thread_id='thread-1')
        found = _find_case(None, 'thread-2', 'A10', f'Re: {self.SUBJECT}')
        self.assertIsNone(found)

    # 오픈 키워드가 없는 고객사 스레드 — 회신 접두어를 신호로 이어붙인다
    THREAD_SUBJECT = '[삼성SDS] A10 receive-buffer 설정 건'

    def test_reply_without_open_keyword_matches_existing_case(self):
        case = make_case(vendor='A10')
        make_email(case, self.THREAD_SUBJECT, thread_id='thread-1')
        found = _find_case(None, 'thread-2', 'A10', f'RE:(6) {self.THREAD_SUBJECT}')
        self.assertEqual(found, case)

    def test_reply_matches_case_started_mid_conversation(self):
        # 스레드가 매번 갈리면 케이스의 첫 메일부터 이미 'RE:(4) …' 회신이다
        case = make_case(vendor='A10')
        make_email(case, f'RE:(4) {self.THREAD_SUBJECT}', thread_id='thread-1')
        make_email(case, f'RE: RE:(4) {self.THREAD_SUBJECT}', thread_id='thread-1')
        found = _find_case(None, 'thread-2', 'A10', f'FW: RE:(12) {self.THREAD_SUBJECT}')
        self.assertEqual(found, case)

    def test_resend_with_remind_tag_matches_existing_case(self):
        # '[Remind] FW: RE:(2) 제목' — 재발송 표시가 접두어 앞에 붙어도 같은 대화
        case = make_case(vendor='A10')
        make_email(case, '[문의] GW Health Check 설정', thread_id='thread-1')
        found = _find_case(None, 'thread-2', 'A10',
                           '[Remind] FW: RE:(2) [문의] GW Health Check 설정')
        self.assertEqual(found, case)

    def test_meaningful_bracket_tag_is_kept(self):
        # 뜻이 담긴 대괄호 태그는 벗기지 않는다 — 고객사가 다르면 다른 케이스
        case = make_case(vendor='A10')
        make_email(case, '[삼성SDS] A10 receive-buffer 설정 건', thread_id='thread-1')
        found = _find_case(None, 'thread-2', 'A10', 'RE: A10 receive-buffer 설정 건')
        self.assertIsNone(found)

    def test_different_vendor_case_number_is_not_merged(self):
        # 같은 증상으로 번호만 새로 딴 재접수 — 번호가 제목에 박혀 있어
        # 정리된 제목이 달라지므로 제목 폴백으로 이어붙지 않는다
        case = make_case(vendor='A10', vendor_case_number='455910')
        make_email(case, 'Case # 455910 [Samsung SDS] Increase TH940 Disk Case',
                   thread_id='thread-1')
        found = _find_case('459173', 'thread-2', 'A10',
                           'Re: Case # 459173 [Samsung SDS] Increase TH940 Disk Case')
        self.assertIsNone(found)

    def test_new_conversation_without_reply_prefix_creates_new_case(self):
        # 접두어도 오픈 키워드도 없는 새 제목은 기존 케이스에 붙지 않는다
        case = make_case(vendor='A10')
        make_email(case, f'RE:(4) {self.THREAD_SUBJECT}', thread_id='thread-1')
        found = _find_case(None, 'thread-2', 'A10', self.THREAD_SUBJECT)
        self.assertIsNone(found)


class FindDuplicateCasesTests(TestCase):
    """제목이 같은데 갈린 케이스 묶음 보고 (find_duplicate_cases 관리 명령)."""

    def _run(self):
        out = StringIO()
        call_command('find_duplicate_cases', stdout=out)
        return out.getvalue()

    def test_split_conversation_is_reported_as_mergeable(self):
        first = make_case(vendor='A10', summary='첫 케이스')
        second = make_case(vendor='A10', summary='갈린 케이스')
        make_email(first, '[삼성SDS] A10 receive-buffer 설정 건', thread_id='t1')
        make_email(second, 'RE:(6) [삼성SDS] A10 receive-buffer 설정 건', thread_id='t2')

        output = self._run()
        self.assertIn(first.case_id, output)
        self.assertIn(f'{second.case_id} [Open]', output)
        self.assertIn('병합 대상', output)

    def test_repeated_notice_is_reported_without_merge_mark(self):
        # 공지성 메일은 제목이 같아도 병합 대상이 아니다 — 보고만 한다
        first = make_case(vendor='Arista')
        second = make_case(vendor='Arista')
        make_email(first, 'New End of Sale email notification', thread_id='t1')
        make_email(second, 'New End of Sale email notification', thread_id='t2')

        self.assertNotIn('병합 대상', self._run())

    def test_reopened_case_with_new_number_is_not_marked_mergeable(self):
        # 번호만 새로 딴 재접수는 회신이어도 제목에 번호가 박혀 있어 병합 대상이 아니다
        first = make_case(vendor='A10', vendor_case_number='455910')
        second = make_case(vendor='A10', vendor_case_number='459173')
        make_email(first, 'Re: Case # 455910 [Samsung SDS] Increase TH940 Disk Case',
                   thread_id='t1')
        make_email(first, '[Samsung SDS] Increase TH940 Disk Case', thread_id='t1')
        # 번호 없는 제목이 겹쳐 같은 묶음으로 잡히지만, 첫 메일 제목엔 새 번호가 박혀 있다
        make_email(second, 'Re: Case # 459173 [Samsung SDS] Increase TH940 Disk Case',
                   thread_id='t2')
        make_email(second, '[Samsung SDS] Increase TH940 Disk Case', thread_id='t2')

        output = self._run()
        self.assertIn(first.case_id, output)      # 묶음으로는 보고하되
        self.assertNotIn('병합 대상', output)      # 자동 병합 대상은 아님

    def test_unrelated_cases_are_not_grouped(self):
        first = make_case(vendor='A10')
        second = make_case(vendor='A10')
        make_email(first, '[문의] GW Health Check 설정', thread_id='t1')
        make_email(second, '[문의] DSR 및 IPIP Tunnel 설정 검토', thread_id='t2')

        self.assertIn('중복 의심 묶음이 없습니다', self._run())


class MergeCasesTests(TestCase):
    """중복 생성된 케이스 병합 (merge_cases 관리 명령)."""

    def setUp(self):
        self.target = make_case(vendor='A10', status='Open', summary='receive-buffer 설정 건',
                                action_steps='[2026-07-31 05:14 발신] 설정값 문의')
        self.other = make_case(vendor='A10', status='Pending', summary='receive-buffer 문의',
                               action_steps='[2026-08-03 03:12 수신] 벤더 회신 도착',
                               device_model='TH1040-F')
        make_email(self.target, 'RE:(4) [삼성SDS] receive-buffer 설정 건')
        self.moved = make_email(self.other, 'RE:(7) [삼성SDS] receive-buffer 설정 건')

    def test_merge_moves_emails_and_deletes_source(self):
        call_command('merge_cases', self.target.case_id, self.other.case_id)

        self.moved.refresh_from_db()
        self.assertEqual(self.moved.case_id, self.target.pk)
        self.assertEqual(self.target.emails.count(), 2)
        self.assertFalse(Case.objects.filter(pk=self.other.pk).exists())

    def test_merge_orders_timeline_and_fills_empty_fields(self):
        call_command('merge_cases', self.target.case_id, self.other.case_id)

        self.target.refresh_from_db()
        steps = self.target.action_steps
        self.assertLess(steps.index('설정값 문의'), steps.index('벤더 회신 도착'))
        self.assertIn('병합했습니다', steps)
        self.assertEqual(self.target.device_model, 'TH1040-F')
        # 상태는 가장 최근 메일을 받은 케이스의 것을 따른다
        self.assertEqual(self.target.status, 'Pending')

    def test_merge_target_is_the_case_that_started_the_conversation(self):
        # 케이스 id 순이 아니라 첫 메일이 이른 쪽이 남는다 (동기화 순서 ≠ 대화 순서)
        older_talk = make_case(vendor='A10', summary='먼저 시작된 대화')
        email = make_email(older_talk, 'RE:(2) [삼성SDS] receive-buffer 설정 건')
        email.received_at = timezone.now() - timedelta(days=3)
        email.save()

        call_command('merge_cases', self.target.case_id, older_talk.case_id)

        older_talk.refresh_from_db()
        self.assertEqual(older_talk.emails.count(), 2)
        self.assertFalse(Case.objects.filter(pk=self.target.pk).exists())

    def test_merge_rejects_different_vendors(self):
        arista = make_case(vendor='Arista')
        with self.assertRaises(CommandError):
            call_command('merge_cases', self.target.case_id, arista.case_id)


class HelpAgentToolTests(TestCase):
    """헬프 에이전트 DB 조회 도구의 동작 검증 (LLM 호출 없음)."""

    def setUp(self):
        self.case = make_case(
            vendor='A10', status='Resolved',
            summary='수원 SCPv2 DATALB 파티션 변경 오류',
            device_model='TH1040-F',
        )
        make_email(self.case, '[Caseopen] 수원 SCPv2 DATALB 파티션 변경 오류')
        make_case(vendor='Arista', summary='40G Interface Link FLAP')

    def test_search_by_keyword_and_vendor(self):
        data = json.loads(help_agent._search_cases(query='파티션', vendor='A10'))
        self.assertEqual(len(data['results']), 1)
        self.assertEqual(data['results'][0]['case_id'], self.case.case_id)

    def test_search_by_case_ref(self):
        data = json.loads(help_agent._search_cases(query=self.case.case_id))
        self.assertEqual(data['results'][0]['case_id'], self.case.case_id)

    def test_detail_resolves_c_format_and_includes_emails(self):
        data = json.loads(help_agent._get_case_detail(self.case.case_id))
        self.assertEqual(data['case_id'], self.case.case_id)
        self.assertEqual(len(data['emails']), 1)

    def test_detail_unknown_case_returns_error(self):
        data = json.loads(help_agent._get_case_detail('C-9999'))
        self.assertIn('error', data)

    def test_stats_counts_by_vendor(self):
        data = json.loads(help_agent._get_case_stats(days=7))
        self.assertEqual(data['total'], 2)
        self.assertEqual(data['by_vendor']['A10'], 1)

    def test_verify_flags_hallucinated_case_ref(self):
        reply = help_agent._verify_case_refs(
            f'{self.case.case_id} 및 C-8888 참조')
        self.assertIn('C-8888', reply)
        self.assertIn('확인되지 않았습니다', reply)

    def test_verify_passes_valid_refs_untouched(self):
        reply = help_agent._verify_case_refs(f'{self.case.case_id} 참조')
        self.assertNotIn('확인되지', reply)

    def test_list_recent_cases_marks_new_and_filters_vendor(self):
        data = json.loads(help_agent._list_recent_cases(days=7, vendor='A10'))
        self.assertEqual(data['count'], 1)
        self.assertEqual(data['results'][0]['case_id'], self.case.case_id)
        self.assertTrue(data['results'][0]['is_new'])


class HelpAgentTriageTests(TestCase):
    """트리아지: 규칙 우선, 애매하면 haiku 분류, 실패 시 search 폴백."""

    def test_report_keyword_skips_llm(self):
        client = MagicMock()
        agent = help_agent._triage(
            client, [{'role': 'user', 'content': '이번 주 케이스 리포트 만들어줘'}])
        self.assertEqual(agent, 'report')
        client.messages.create.assert_not_called()

    def test_ambiguous_question_uses_llm_classifier(self):
        client = MagicMock()
        client.messages.create.return_value = SimpleNamespace(
            content=[_fake_block(type='text', text='report')])
        agent = help_agent._triage(
            client, [{'role': 'user', 'content': '요즘 케이스들 어떻게 돌아가?'}])
        self.assertEqual(agent, 'report')
        client.messages.create.assert_called_once()

    def test_classifier_failure_falls_back_to_search(self):
        client = MagicMock()
        client.messages.create.side_effect = anthropic.APIConnectionError(
            request=MagicMock())
        agent = help_agent._triage(
            client, [{'role': 'user', 'content': 'C-1122 상태 알려줘'}])
        self.assertEqual(agent, 'search')

    def test_off_topic_classification(self):
        client = MagicMock()
        client.messages.create.return_value = SimpleNamespace(
            content=[_fake_block(type='text', text='off_topic')])
        agent = help_agent._triage(
            client, [{'role': 'user', 'content': '오늘 저녁 뭐 먹을까?'}])
        self.assertEqual(agent, 'off_topic')

    def test_followup_question_includes_conversation_context(self):
        # "인터넷에서 더 찾아줘" 같은 후속 질문은 직전 맥락과 함께 분류돼야 함
        client = MagicMock()
        client.messages.create.return_value = SimpleNamespace(
            content=[_fake_block(type='text', text='tech')])
        agent = help_agent._triage(client, [
            {'role': 'user', 'content': 'C-1122 VRRP 버그 상태 알려줘'},
            {'role': 'assistant', 'content': 'C-1122는 Resolved 상태입니다.'},
            {'role': 'user', 'content': '인터넷에서 상세 검색해줘'},
        ])
        self.assertEqual(agent, 'tech')
        sent = client.messages.create.call_args.kwargs['messages'][0]['content']
        self.assertIn('이전 대화 맥락', sent)
        self.assertIn('VRRP 버그', sent)
        self.assertIn('인터넷에서 상세 검색해줘', sent)

    @override_settings(ANTHROPIC_API_KEY='test-key', HELP_AGENT_MODEL='claude-haiku-4-5')
    def test_off_topic_short_circuits_without_agent_call(self):
        fake_client = MagicMock()
        fake_client.messages.create.return_value = SimpleNamespace(
            content=[_fake_block(type='text', text='off_topic')])
        with patch.object(help_agent.anthropic, 'Anthropic', return_value=fake_client):
            result = help_agent.chat([{'role': 'user', 'content': '주식 추천해줘'}])

        self.assertEqual(result['agent'], 'off_topic')
        self.assertEqual(result['tool_calls'], [])
        self.assertIn('범위 밖', result['reply'])
        # 트리아지 1회만 호출 — 에이전트 본 호출 없음 (비용 가드)
        self.assertEqual(fake_client.messages.create.call_count, 1)


def _fake_block(**kwargs):
    return SimpleNamespace(**kwargs)


class HelpAgentChatLoopTests(TestCase):
    """에이전트 루프: 도구 호출 → 결과 회신 → 최종 답변 (Anthropic 모킹)."""

    def setUp(self):
        self.case = make_case(vendor='A10', summary='VRRP failover 장애')

    @override_settings(ANTHROPIC_API_KEY='test-key', HELP_AGENT_MODEL='claude-haiku-4-5')
    def test_tool_loop_returns_final_reply_and_trace(self):
        triage_turn = SimpleNamespace(
            stop_reason='end_turn',
            content=[_fake_block(type='text', text='search')],
        )
        tool_turn = SimpleNamespace(
            stop_reason='tool_use',
            content=[_fake_block(type='tool_use', id='tu_1', name='search_cases',
                                 input={'query': 'VRRP'})],
        )
        final_turn = SimpleNamespace(
            stop_reason='end_turn',
            content=[_fake_block(type='text',
                                 text=f'{self.case.case_id} 케이스가 있습니다.')],
        )
        fake_client = MagicMock()
        fake_client.messages.create.side_effect = [triage_turn, tool_turn, final_turn]

        with patch.object(help_agent.anthropic, 'Anthropic', return_value=fake_client):
            result = help_agent.chat([{'role': 'user', 'content': 'VRRP 장애 사례 찾아줘'}])

        self.assertEqual(result['agent'], 'search')
        self.assertIn(self.case.case_id, result['reply'])
        self.assertEqual(result['tool_calls'], [{'name': 'search_cases',
                                                 'input': {'query': 'VRRP'}}])
        # 3번째 호출(도구 회신 후)의 messages에 tool_result가 포함됐는지 확인
        third_call_messages = fake_client.messages.create.call_args_list[2].kwargs['messages']
        self.assertEqual(third_call_messages[-1]['content'][0]['type'], 'tool_result')

    @override_settings(ANTHROPIC_API_KEY='test-key',
                       HELP_AGENT_MODEL='claude-haiku-4-5',
                       REPORT_AGENT_MODEL='claude-sonnet-5')
    def test_report_request_routes_to_report_agent(self):
        final_turn = SimpleNamespace(
            stop_reason='end_turn',
            content=[_fake_block(type='text', text='## 주간 리포트\n요약입니다.')],
        )
        fake_client = MagicMock()
        # 리포팅은 문서 스킬 때문에 beta 엔드포인트를 사용한다
        fake_client.beta.messages.create.return_value = final_turn

        with patch.object(help_agent.anthropic, 'Anthropic', return_value=fake_client):
            result = help_agent.chat(
                [{'role': 'user', 'content': '이번 주 케이스 리포트 작성해줘'}])

        self.assertEqual(result['agent'], 'report')
        self.assertEqual(result['model'], 'claude-sonnet-5')
        self.assertNotIn('files', result)  # 문서를 안 만들면 files 없음
        # 리포트 키워드는 규칙 분기 → 트리아지 LLM 호출 없이 본 호출 1회만
        fake_client.messages.create.assert_not_called()
        call = fake_client.beta.messages.create.call_args_list[0]
        self.assertEqual(call.kwargs['model'], 'claude-sonnet-5')
        tool_names = [t['name'] for t in call.kwargs['tools']]
        self.assertIn('list_recent_cases', tool_names)
        # 문서 스킬 구성: code_execution 도구 + 스킬 컨테이너 + beta 헤더
        self.assertIn('code_execution', tool_names)
        skill_ids = [s['skill_id'] for s in call.kwargs['container']['skills']]
        self.assertEqual(skill_ids, ['docx', 'xlsx', 'pptx'])
        self.assertEqual(call.kwargs['betas'],
                         ['code-execution-2025-08-25', 'skills-2025-10-02'])

    @override_settings(ANTHROPIC_API_KEY='test-key',
                       HELP_AGENT_MODEL='claude-haiku-4-5',
                       REPORT_AGENT_MODEL='claude-sonnet-5')
    def test_report_collects_generated_files(self):
        # 코드 실행 결과 블록 안에 중첩된 file_id를 수집하는지
        file_ref = _fake_block(type='bash_code_execution_output',
                               file_id='file_abc123')
        exec_result = _fake_block(
            type='bash_code_execution_tool_result',
            content=_fake_block(type='bash_code_execution_result',
                                content=[file_ref]),
        )
        final_turn = SimpleNamespace(
            stop_reason='end_turn',
            content=[exec_result,
                     _fake_block(type='text', text='엑셀 리포트를 만들었습니다.')],
        )
        fake_client = MagicMock()
        fake_client.beta.messages.create.return_value = final_turn
        fake_client.beta.files.retrieve_metadata.return_value = SimpleNamespace(
            filename='caseflow_report.xlsx', size_bytes=2048)

        with patch.object(help_agent.anthropic, 'Anthropic', return_value=fake_client):
            result = help_agent.chat(
                [{'role': 'user', 'content': '이번 주 리포트를 엑셀로 작성해줘'}])

        self.assertEqual(result['files'], [{
            'file_id': 'file_abc123',
            'filename': 'caseflow_report.xlsx',
            'size_bytes': 2048,
        }])

    def test_describe_files_filters_non_documents_and_dedupes(self):
        fake_client = MagicMock()
        fake_client.beta.files.retrieve_metadata.side_effect = [
            SimpleNamespace(filename='report.docx', size_bytes=100),
            SimpleNamespace(filename='build_report.py', size_bytes=50),
        ]
        files = help_agent._describe_files(
            fake_client, ['file_doc', 'file_doc', 'file_script'])
        self.assertEqual([f['filename'] for f in files], ['report.docx'])
        # 중복 file_id는 메타데이터 조회도 1회만
        self.assertEqual(fake_client.beta.files.retrieve_metadata.call_count, 2)

    @override_settings(ANTHROPIC_API_KEY='')
    def test_missing_api_key_raises(self):
        with self.assertRaises(RuntimeError):
            help_agent.chat([{'role': 'user', 'content': '안녕'}])


class HelpAgentTemplateTests(TestCase):
    """사내 템플릿 모드 — 워딩 트리거, 파일 첨부, 해시 캐싱 (Anthropic 모킹)."""

    def setUp(self):
        # 실제 템플릿 파일은 gitignore 대상이라 테스트는 임시 파일로 대체한다
        self.tmpdir = tempfile.TemporaryDirectory()
        self.addCleanup(self.tmpdir.cleanup)
        docx = Path(self.tmpdir.name) / 'demo.docx'
        docx.write_bytes(b'docx-template-v1')
        self.docx_patch = patch.dict(
            help_agent.REPORT_TEMPLATES['docx'], {'path': docx})
        self.docx_patch.start()
        self.addCleanup(self.docx_patch.stop)

    def test_match_template_requires_both_keywords(self):
        # '템플릿' + 형식 단어가 함께 있을 때만 반응 (일반 파일 생성과 구분)
        self.assertEqual(
            help_agent._match_template('사내보고서 워드 템플릿으로 작성해줘'), 'docx')
        self.assertEqual(
            help_agent._match_template('C-1122 PPT 템플릿으로 정리해줘'), 'pptx')
        self.assertIsNone(help_agent._match_template('이번 주 리포트를 워드로 작성해줘'))
        self.assertIsNone(help_agent._match_template('템플릿이 뭐야?'))

    def test_template_wording_routes_to_report_without_llm_triage(self):
        # "…템플릿으로 만들어줘"는 보고서 단어가 없어도 규칙 분기로 report에 가야 함
        self.assertIn('템플릿', help_agent.REPORT_KEYWORDS)

    @override_settings(ANTHROPIC_API_KEY='test-key',
                       HELP_AGENT_MODEL='claude-haiku-4-5',
                       REPORT_AGENT_MODEL='claude-sonnet-5')
    def test_template_request_attaches_file_and_addendum(self):
        final_turn = SimpleNamespace(
            stop_reason='end_turn',
            content=[_fake_block(type='text', text='템플릿 보고서를 만들었습니다.')])
        fake_client = MagicMock()
        fake_client.beta.messages.create.return_value = final_turn
        fake_client.beta.files.upload.return_value = SimpleNamespace(id='file_tpl_1')

        with patch.object(help_agent.anthropic, 'Anthropic', return_value=fake_client):
            help_agent.chat(
                [{'role': 'user', 'content': 'C-1122 사내보고서 워드 템플릿으로 작성해줘'}])

        call = fake_client.beta.messages.create.call_args_list[0]
        last_content = call.kwargs['messages'][-1]['content']
        self.assertEqual(last_content[0],
                         {'type': 'container_upload', 'file_id': 'file_tpl_1'})
        self.assertIn('워드 템플릿으로 작성해줘', last_content[1]['text'])
        self.assertIn('사내 템플릿 모드', call.kwargs['system'])
        # 해시:file_id 캐시 저장 확인
        self.assertTrue(AppSetting.get('report_template_docx').endswith(':file_tpl_1'))

    @override_settings(ANTHROPIC_API_KEY='test-key',
                       HELP_AGENT_MODEL='claude-haiku-4-5',
                       REPORT_AGENT_MODEL='claude-sonnet-5')
    def test_plain_report_request_does_not_attach_template(self):
        final_turn = SimpleNamespace(
            stop_reason='end_turn',
            content=[_fake_block(type='text', text='## 주간 리포트')])
        fake_client = MagicMock()
        fake_client.beta.messages.create.return_value = final_turn

        with patch.object(help_agent.anthropic, 'Anthropic', return_value=fake_client):
            help_agent.chat([{'role': 'user', 'content': '이번 주 리포트 작성해줘'}])

        fake_client.beta.files.upload.assert_not_called()
        call = fake_client.beta.messages.create.call_args_list[0]
        self.assertIsInstance(call.kwargs['messages'][-1]['content'], str)
        self.assertNotIn('사내 템플릿 모드', call.kwargs['system'])

    def test_file_id_cached_until_template_file_changes(self):
        fake_client = MagicMock()
        fake_client.beta.files.upload.side_effect = [
            SimpleNamespace(id='file_v1'), SimpleNamespace(id='file_v2')]

        self.assertEqual(help_agent._template_file_id(fake_client, 'docx'), 'file_v1')
        # 같은 파일이면 재업로드 없이 캐시 사용
        self.assertEqual(help_agent._template_file_id(fake_client, 'docx'), 'file_v1')
        self.assertEqual(fake_client.beta.files.upload.call_count, 1)

        # 파일 교체(해시 변경) → 재업로드 + 옛 파일 삭제
        help_agent.REPORT_TEMPLATES['docx']['path'].write_bytes(b'docx-template-v2')
        self.assertEqual(help_agent._template_file_id(fake_client, 'docx'), 'file_v2')
        fake_client.beta.files.delete.assert_called_once_with('file_v1')

    @override_settings(ANTHROPIC_API_KEY='test-key',
                       HELP_AGENT_MODEL='claude-haiku-4-5',
                       REPORT_AGENT_MODEL='claude-sonnet-5')
    def test_upload_failure_falls_back_to_plain_report(self):
        # 템플릿 첨부 실패는 500 대신 일반 리포트로 진행 (시연 중단 방지)
        final_turn = SimpleNamespace(
            stop_reason='end_turn',
            content=[_fake_block(type='text', text='일반 보고서입니다.')])
        fake_client = MagicMock()
        fake_client.beta.messages.create.return_value = final_turn
        fake_client.beta.files.upload.side_effect = anthropic.APIConnectionError(
            request=MagicMock())

        with patch.object(help_agent.anthropic, 'Anthropic', return_value=fake_client):
            result = help_agent.chat(
                [{'role': 'user', 'content': '워드 템플릿으로 보고서 작성해줘'}])

        self.assertEqual(result['reply'], '일반 보고서입니다.')
        call = fake_client.beta.messages.create.call_args_list[0]
        self.assertNotIn('사내 템플릿 모드', call.kwargs['system'])


class HelpAgentEndpointTests(TestCase):
    """POST /api/help-agent/chat/ 의 인증·검증·응답.

    엔지니어 이상 사용 가능 (2026-07-21, 관리자 전용에서 확대).
    """

    def setUp(self):
        from .permissions import set_user_role
        viewer = User.objects.create_user('viewer1', password='pw123456')
        set_user_role(viewer, 'viewer')
        User.objects.create_user('admin1', password='pw123456', is_staff=True)

    def login(self, username):
        self.client.post('/api/auth/login/',
                         {'username': username, 'password': 'pw123456'},
                         content_type='application/json')

    def test_requires_login(self):
        res = self.client.post('/api/help-agent/chat/',
                               {'messages': [{'role': 'user', 'content': '안녕'}]},
                               content_type='application/json')
        self.assertIn(res.status_code, (401, 403))

    def test_viewer_is_blocked(self):
        self.login('viewer1')
        res = self.client.post('/api/help-agent/chat/',
                               {'messages': [{'role': 'user', 'content': '안녕'}]},
                               content_type='application/json')
        self.assertEqual(res.status_code, 403)

    def test_invalid_payload_rejected(self):
        self.login('admin1')
        res = self.client.post('/api/help-agent/chat/', {'messages': []},
                               content_type='application/json')
        self.assertEqual(res.status_code, 400)
        res = self.client.post(
            '/api/help-agent/chat/',
            {'messages': [{'role': 'assistant', 'content': '내가 마지막'}]},
            content_type='application/json')
        self.assertEqual(res.status_code, 400)

    def test_admin_can_chat(self):
        self.login('admin1')
        with patch('api.views.help_agent.chat',
                   return_value={'reply': '안녕하세요', 'tool_calls': [], 'model': 'm'}):
            res = self.client.post(
                '/api/help-agent/chat/',
                {'messages': [{'role': 'user', 'content': '안녕'}]},
                content_type='application/json')
        self.assertEqual(res.status_code, 200)
        self.assertEqual(res.json()['reply'], '안녕하세요')

    def test_file_download_blocks_viewer(self):
        self.login('viewer1')
        res = self.client.get('/api/help-agent/files/file_abc123/')
        self.assertEqual(res.status_code, 403)

    def test_file_download_rejects_invalid_id(self):
        self.login('admin1')
        res = self.client.get('/api/help-agent/files/not-a-file-id/')
        self.assertEqual(res.status_code, 400)

    def test_file_download_relays_content(self):
        self.login('admin1')
        with patch('api.views.help_agent.download_file',
                   return_value=(
                       '리포트.xlsx',
                       'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                       b'excel-bytes')):
            res = self.client.get('/api/help-agent/files/file_abc123/')
        self.assertEqual(res.status_code, 200)
        self.assertEqual(res.content, b'excel-bytes')
        self.assertIn("filename*=UTF-8''%EB%A6%AC%ED%8F%AC%ED%8A%B8.xlsx",
                      res['Content-Disposition'])


@override_settings(SEARCH_BLOCKED_TERMS=['samsung', '삼성', '하나은행'])
class SearchQuerySanitizerTests(TestCase):
    """웹 검색어 보안 정제 — 고객사명·시리얼·사설 IP 제거 (코드 가드레일)."""

    def test_customer_names_removed(self):
        clean, removed = help_agent._sanitize_search_query(
            '삼성 SCP 환경 A10 파티션 변경 오류')
        self.assertNotIn('삼성', clean)
        self.assertIn('A10 파티션 변경 오류', clean)
        self.assertIn('삼성', removed)

    def test_private_ip_and_serial_removed(self):
        clean, removed = help_agent._sanitize_search_query(
            'TH1040-F 10.20.30.40 TH10154022070160 failover 원인')
        self.assertNotIn('10.20.30.40', clean)
        self.assertNotIn('TH10154022070160', clean)
        self.assertIn('TH1040-F', clean)  # 모델명은 유지
        self.assertIn('failover', clean)
        self.assertEqual(len(removed), 2)

    def test_clean_technical_query_untouched(self):
        clean, removed = help_agent._sanitize_search_query(
            'Arista EOS 4.32.4M PhyEthtool log advisory')
        self.assertEqual(clean, 'Arista EOS 4.32.4M PhyEthtool log advisory')
        self.assertEqual(removed, [])

    def test_vendor_bug_id_is_kept(self):
        # ACOS-104904 같은 버그 ID는 시리얼이 아니다 — 검색에 필요 (오탐 회귀 방지)
        clean, removed = help_agent._sanitize_search_query(
            'A10 ACOS-104904 VRRP-A advertisement timer bug')
        self.assertIn('ACOS-104904', clean)
        self.assertEqual(removed, [])


class WebSearchToolTests(TestCase):
    """web_search 도구 — Serper 연동(모킹)과 키 미설정 처리."""

    @override_settings(SERPER_API_KEY='')
    def test_missing_key_returns_error(self):
        data = json.loads(help_agent._web_search('EOS bug'))
        self.assertIn('error', data)

    @override_settings(SERPER_API_KEY='k', SEARCH_BLOCKED_TERMS=['삼성'])
    def test_results_parsed_and_sanitize_notice(self):
        fake_response = MagicMock()
        fake_response.json.return_value = {'organic': [
            {'title': 'ACOS Release Notes', 'link': 'https://a10.com/rn',
             'snippet': '6.0.9 fixes'},
        ]}
        with patch.object(help_agent.httpx, 'post',
                          return_value=fake_response) as post:
            data = json.loads(help_agent._web_search('삼성 ACOS 6.0.8 bug'))

        self.assertEqual(data['results'][0]['url'], 'https://a10.com/rn')
        self.assertIn('제거됨', data['notice'])
        # 실제 전송된 검색어에 고객사명이 없어야 함
        sent_query = post.call_args.kwargs['json']['q']
        self.assertNotIn('삼성', sent_query)

    @override_settings(SERPER_API_KEY='k', SEARCH_BLOCKED_TERMS=['삼성'])
    def test_fully_blocked_query_returns_error_without_request(self):
        with patch.object(help_agent.httpx, 'post') as post:
            data = json.loads(help_agent._web_search('삼성'))
        self.assertIn('error', data)
        post.assert_not_called()


class TechAgentFlowTests(TestCase):
    """② 기술지원: 트리아지 → sonnet 답변 → haiku 검수 → (미흡 시) 수정."""

    def _triage_resp(self, label):
        return SimpleNamespace(stop_reason='end_turn',
                               content=[_fake_block(type='text', text=label)])

    def _text_resp(self, text, stop_reason='end_turn'):
        return SimpleNamespace(stop_reason=stop_reason,
                               content=[_fake_block(type='text', text=text)])

    @override_settings(ANTHROPIC_API_KEY='test-key',
                       HELP_AGENT_MODEL='claude-haiku-4-5',
                       TECH_AGENT_MODEL='claude-sonnet-5')
    def test_evaluator_pass_returns_reply_without_revision(self):
        fake_client = MagicMock()
        fake_client.messages.create.side_effect = [
            self._triage_resp('tech'),
            self._text_resp('ACOS 6.0.9에서 수정되었습니다. [RN](https://a10.com/rn)'),
            self._text_resp('{"ok": true}'),  # 평가자
        ]
        with patch.object(help_agent.anthropic, 'Anthropic', return_value=fake_client):
            result = help_agent.chat(
                [{'role': 'user', 'content': 'ACOS 6.0.8 VRRP 버그 수정 버전 알려줘'}])

        self.assertEqual(result['agent'], 'tech')
        self.assertEqual(result['model'], 'claude-sonnet-5')
        self.assertTrue(result['evaluation']['ok'])
        self.assertIn('6.0.9', result['reply'])
        self.assertEqual(fake_client.messages.create.call_count, 3)  # 수정 라운드 없음

    @override_settings(ANTHROPIC_API_KEY='test-key',
                       HELP_AGENT_MODEL='claude-haiku-4-5',
                       TECH_AGENT_MODEL='claude-sonnet-5')
    def test_evaluator_fail_triggers_one_revision(self):
        fake_client = MagicMock()
        fake_client.messages.create.side_effect = [
            self._triage_resp('tech'),
            self._text_resp('근거 없는 초안'),
            self._text_resp('{"ok": false, "issues": ["출처 인용 없음"]}'),  # 평가자
            self._text_resp('수정된 답변 [출처](https://vendor.com/doc)'),   # 수정 라운드
        ]
        with patch.object(help_agent.anthropic, 'Anthropic', return_value=fake_client):
            result = help_agent.chat(
                [{'role': 'user', 'content': 'EOS 업그레이드 시 주의사항 알려줘'}])

        self.assertFalse(result['evaluation']['ok'])
        self.assertIn('수정된 답변', result['reply'])
        self.assertEqual(fake_client.messages.create.call_count, 4)
        # 수정 요청에 검수 피드백이 전달됐는지 확인
        revision_messages = fake_client.messages.create.call_args_list[3].kwargs['messages']
        self.assertIn('자동 검수 피드백', revision_messages[-1]['content'])

    @override_settings(ANTHROPIC_API_KEY='test-key',
                       HELP_AGENT_MODEL='claude-haiku-4-5',
                       TECH_AGENT_MODEL='claude-sonnet-5')
    def test_revision_preamble_does_not_replace_the_answer(self):
        """수정 라운드가 "다시 확인하겠습니다" 예고문을 내놓으면 원본을 지킨다.

        운영에서 후속 질문마다 40자짜리 예고문만 저장돼 답변이 사라진 것처럼
        보였던 장애의 회귀 테스트.
        """
        answer = '## TACACS+ 설정\n' + 'ACOS 설정 절차 상세 설명. ' * 40
        fake_client = MagicMock()
        fake_client.messages.create.side_effect = [
            self._triage_resp('tech'),
            self._text_resp(answer),
            self._text_resp('{"ok": false, "issues": ["출처 인용 없음"]}'),
            self._text_resp('피드백을 반영하여 문서를 직접 확인하겠습니다.'),
        ]
        with patch.object(help_agent.anthropic, 'Anthropic', return_value=fake_client):
            result = help_agent.chat(
                [{'role': 'user', 'content': 'A10 TACACS 연동 가이드 알려줘'},
                 {'role': 'assistant', 'content': answer},
                 {'role': 'user', 'content': '진행해줘'}])

        self.assertIn('TACACS+ 설정', result['reply'])
        self.assertNotIn('확인하겠습니다', result['reply'])

    @override_settings(ANTHROPIC_API_KEY='test-key',
                       HELP_AGENT_MODEL='claude-haiku-4-5',
                       TECH_AGENT_MODEL='claude-sonnet-5')
    def test_revision_round_blocks_tool_use(self):
        """수정 라운드는 도구를 못 쓰게 막아 예고문 유도 자체를 차단한다."""
        fake_client = MagicMock()
        fake_client.messages.create.side_effect = [
            self._triage_resp('tech'),
            self._text_resp('근거 없는 초안'),
            self._text_resp('{"ok": false, "issues": ["출처 인용 없음"]}'),
            self._text_resp('수정된 답변 [출처](https://vendor.com/doc)'),
        ]
        with patch.object(help_agent.anthropic, 'Anthropic', return_value=fake_client):
            help_agent.chat([{'role': 'user', 'content': 'EOS 주의사항 알려줘'}])

        revision_kwargs = fake_client.messages.create.call_args_list[3].kwargs
        self.assertEqual(revision_kwargs['tool_choice'], {'type': 'none'})

    @override_settings(ANTHROPIC_API_KEY='test-key',
                       HELP_AGENT_MODEL='claude-haiku-4-5',
                       TECH_AGENT_MODEL='claude-sonnet-5')
    def test_evaluator_receives_previous_context(self):
        """짧은 후속 지시("진행해줘")도 검수자가 주제를 알 수 있어야 한다."""
        fake_client = MagicMock()
        fake_client.messages.create.side_effect = [
            self._triage_resp('tech'),
            self._text_resp('A10 TACACS 설정 [출처](https://a10.com/doc)'),
            self._text_resp('{"ok": true}'),
        ]
        with patch.object(help_agent.anthropic, 'Anthropic', return_value=fake_client):
            help_agent.chat(
                [{'role': 'user', 'content': 'A10 TACACS 연동 가이드 알려줘'},
                 {'role': 'assistant', 'content': 'TACACS+ 서버 등록 절차입니다'},
                 {'role': 'user', 'content': '진행해줘'}])

        evaluator_content = (fake_client.messages.create
                             .call_args_list[2].kwargs['messages'][0]['content'])
        self.assertIn('이전 대화 맥락', evaluator_content)
        self.assertIn('TACACS 연동 가이드', evaluator_content)
        self.assertIn('[현재 질문]\n진행해줘', evaluator_content)

    @override_settings(ANTHROPIC_API_KEY='test-key',
                       HELP_AGENT_MODEL='claude-haiku-4-5',
                       TECH_AGENT_MODEL='claude-sonnet-5')
    def test_handoff_reroutes_to_target_agent_once(self):
        # 검색 에이전트에게 웹 검색 요청이 잘못 배정 → [HANDOFF:tech] → 재배정
        fake_client = MagicMock()
        fake_client.messages.create.side_effect = [
            self._triage_resp('search'),                 # 트리아지 오분류
            self._text_resp('[HANDOFF:tech]'),           # 검색 에이전트가 핸드오프
            self._text_resp('EOS 4.32 관련 자료입니다. [출처](https://arista.com)'),
            self._text_resp('{"ok": true}'),             # 평가자
        ]
        with patch.object(help_agent.anthropic, 'Anthropic', return_value=fake_client):
            result = help_agent.chat(
                [{'role': 'user', 'content': '인터넷에서 상세 검색해줘'}])

        self.assertEqual(result['agent'], 'tech')
        self.assertEqual(result['model'], 'claude-sonnet-5')
        self.assertIn('EOS 4.32', result['reply'])
        self.assertNotIn('HANDOFF', result['reply'])
        # 재배정된 tech 호출이 tech 프롬프트로 나갔는지 확인
        tech_call = fake_client.messages.create.call_args_list[2]
        self.assertEqual(tech_call.kwargs['model'], 'claude-sonnet-5')

    @override_settings(ANTHROPIC_API_KEY='test-key',
                       HELP_AGENT_MODEL='claude-haiku-4-5',
                       TECH_AGENT_MODEL='claude-sonnet-5')
    def test_handoff_to_same_agent_is_ignored(self):
        # 자기 자신으로의 핸드오프는 재배정하지 않고 마커만 제거 (루프 방지)
        fake_client = MagicMock()
        fake_client.messages.create.side_effect = [
            self._triage_resp('search'),
            self._text_resp('[HANDOFF:search]'),
        ]
        with patch.object(help_agent.anthropic, 'Anthropic', return_value=fake_client):
            result = help_agent.chat([{'role': 'user', 'content': '케이스 찾아줘'}])

        self.assertEqual(result['agent'], 'search')
        self.assertNotIn('HANDOFF', result['reply'])
        self.assertEqual(fake_client.messages.create.call_count, 2)


class GmailSyncConcurrencyTests(TestCase):
    """동기화 동시 실행 잠금과 저장 직전 중복(경쟁 상태) 방어."""

    def test_concurrent_sync_rejected_by_lock(self):
        import fcntl
        from .services import gmail_sync

        # 다른 동기화가 실행 중인 상태를 재현: 잠금을 직접 잡아둔다
        holder = open(gmail_sync._LOCK_FILE, 'w')
        fcntl.flock(holder, fcntl.LOCK_EX | fcntl.LOCK_NB)
        try:
            with self.assertRaises(gmail_sync.SyncInProgress):
                gmail_sync.sync_gmail()
        finally:
            fcntl.flock(holder, fcntl.LOCK_UN)
            holder.close()

        # 잠금 해제 후에는 정상 진입 (Gmail 호출은 mock)
        with patch.object(gmail_sync, '_sync_gmail', return_value={'fetched': 0}):
            self.assertEqual(gmail_sync.sync_gmail(), {'fetched': 0})

    def test_duplicate_at_save_time_rolls_back_new_case(self):
        """중복 체크 통과 후 다른 동기화가 먼저 저장한 경우:
        skipped 처리되고, 이 실행이 만들던 새 케이스도 롤백되어야 한다."""
        from .services import gmail_sync

        other_case = make_case(vendor='Arista')

        message = {
            'id': 'race-msg-1',
            'threadId': 'race-thread-1',
            'internalDate': '0',
            'payload': {'headers': [
                {'name': 'From', 'value': 'Arista Support <support@arista.com>'},
                {'name': 'To', 'value': 'adc@ubersys.co.kr'},
                {'name': 'Subject', 'value': 'New Case: SR 77001 something broken'},
                {'name': 'Date', 'value': 'Mon, 13 Jul 2026 10:00:00 +0900'},
            ]},
        }

        # AI 분석이 도는 사이에 다른 동기화가 같은 메일을 먼저 저장하는 상황
        def analyze_and_race(**kwargs):
            make_email(other_case, 'raced', message_id='race-msg-1',
                       thread_id='race-thread-1')
            return None

        cases_before = Case.objects.count()
        with patch.object(gmail_sync, 'analyze_email', side_effect=analyze_and_race):
            result = gmail_sync._process_message(message)

        self.assertEqual(result, 'skipped')
        # 이메일은 먼저 저장된 1건만 존재
        self.assertEqual(
            CaseEmail.objects.filter(gmail_message_id='race-msg-1').count(), 1)
        # 이 실행이 만들던 새 케이스는 롤백되어 빈 케이스가 남지 않는다
        self.assertEqual(Case.objects.count(), cases_before)


class UsageEventTests(TestCase):
    """파일럿 사용 로그 — 기록 훅과 통계 API."""

    def setUp(self):
        from .permissions import set_user_role
        for username, role in (('uv1', 'viewer'), ('ua1', 'admin')):
            user = User.objects.create_user(username, password='usage-pass-123!')
            set_user_role(user, role)
        self.case = make_case(vendor='A10', summary='사용 로그 테스트 케이스')

    def login(self, username):
        self.client.post('/api/auth/login/',
                         {'username': username, 'password': 'usage-pass-123!'},
                         content_type='application/json')

    def test_login_and_case_views_are_logged(self):
        from .models import UsageEvent
        self.login('uv1')
        self.client.get('/api/cases/')
        self.client.get(f'/api/cases/{self.case.id}/')
        events = list(UsageEvent.objects.order_by('id').values_list('event', flat=True))
        self.assertEqual(events, ['login', 'case_list', 'case_view'])
        detail_event = UsageEvent.objects.get(event='case_view')
        self.assertEqual(detail_event.detail, self.case.case_id)
        self.assertEqual(detail_event.user.username, 'uv1')

    def test_client_search_event_whitelist(self):
        from .models import UsageEvent
        self.login('uv1')
        ok = self.client.post('/api/usage/', {'event': 'search', 'detail': 'VRRP'},
                              content_type='application/json')
        self.assertEqual(ok.status_code, 201)
        self.assertTrue(UsageEvent.objects.filter(event='search', detail='VRRP').exists())
        # 허용 목록 밖 이벤트는 거부 — 지표 오염 방지
        bad = self.client.post('/api/usage/', {'event': 'agent_chat', 'detail': 'x'},
                               content_type='application/json')
        self.assertEqual(bad.status_code, 400)

    def test_stats_admin_only_and_aggregates(self):
        self.login('uv1')
        self.client.get('/api/cases/')
        self.client.post('/api/usage/', {'event': 'search', 'detail': 'failover'},
                         content_type='application/json')
        self.assertEqual(self.client.get('/api/usage/stats/').status_code, 403)

        self.login('ua1')
        res = self.client.get('/api/usage/stats/')
        self.assertEqual(res.status_code, 200)
        data = res.json()
        self.assertEqual(data['active_users'], 2)  # uv1 + ua1(로그인 이벤트)
        self.assertEqual(data['by_event']['search'], 1)
        self.assertEqual(data['by_event']['case_list'], 1)
        usernames = {u['username'] for u in data['users']}
        self.assertEqual(usernames, {'uv1', 'ua1'})
        uv1 = next(u for u in data['users'] if u['username'] == 'uv1')
        self.assertEqual(uv1['searches'], 1)
        self.assertEqual(uv1['logins'], 1)


class KnowledgeFieldsTests(TestCase):
    """지식 본문 8칸 — AI가 채우고, 못 채운 칸은 엔지니어가 화면에서 채운다."""

    FULL = {
        'has_knowledge': True, 'title': 'VRRP failover 지연',
        'environment': 'HA 구성, TH1040-F, ACOS 6.0.8',
        'problem': 'failover가 3초 이상 지연됩니다.',
        'diagnosis': 'show vrrp-a detail 로 advertisement 주기를 확인했습니다.',
        'root_cause': 'advertisement interval 기본값이 큽니다.',
        'resolution': 'vrrp-a interval 100\nwrite memory',
        'verification': 'show vrrp-a detail 에서 interval 100 확인',
        'caveats': '인터벌을 너무 줄이면 플래핑이 발생할 수 있습니다.',
        'related_refs': 'ACOS-104904\nC-1118',
        'device_model': 'TH1040-F', 'software_version': '6.0.8',
    }

    def setUp(self):
        from .permissions import set_user_role
        for username, role in (('kf-v', 'viewer'), ('kf-e', 'engineer')):
            user = User.objects.create_user(username, password='kf-pass-1!')
            set_user_role(user, role)
        self.case = make_case(vendor='A10', status='Resolved', resolution='조치')

    def login(self, username):
        self.client.post('/api/auth/login/',
                         {'username': username, 'password': 'kf-pass-1!'},
                         content_type='application/json')

    def extract(self, payload):
        from .services import knowledge
        with patch.object(knowledge, 'generate_structured_with_model',
                          return_value=(MODEL, payload)), \
             patch.object(knowledge, 'enrich_with_references'):
            return knowledge.extract_knowledge(self.case)

    def test_all_eight_fields_are_saved(self):
        outcome, item = self.extract(dict(self.FULL))
        self.assertEqual(outcome, 'created')
        for field in ('environment', 'diagnosis', 'verification', 'caveats', 'related_refs'):
            self.assertEqual(getattr(item, field), self.FULL[field], field)

    def test_missing_fields_become_blank_not_crash(self):
        """예전 스키마 응답이나 일부 필드를 뺀 응답도 저장은 되어야 한다."""
        outcome, item = self.extract({
            'has_knowledge': True, 'title': '제목', 'problem': '문제',
            'root_cause': '', 'resolution': '조치',
            'device_model': '', 'software_version': '',
        })
        self.assertEqual(outcome, 'created')
        self.assertEqual(item.environment, '')
        self.assertEqual(item.caveats, '')

    def test_engineer_can_fill_blank_fields(self):
        """AI가 비워둔 칸을 사람이 채우는 것이 이 필드들의 목적이다."""
        _, item = self.extract({
            'has_knowledge': True, 'title': '제목', 'problem': '문제',
            'root_cause': '', 'resolution': '조치',
            'device_model': '', 'software_version': '',
        })
        self.login('kf-e')
        res = self.client.patch(f'/api/knowledge/{item.id}/',
                                {'caveats': '재부팅이 필요합니다.',
                                 'verification': 'show version 으로 확인'},
                                content_type='application/json')
        self.assertEqual(res.status_code, 200)
        item.refresh_from_db()
        self.assertEqual(item.caveats, '재부팅이 필요합니다.')
        self.assertEqual(item.verification, 'show version 으로 확인')

    def test_detail_api_exposes_new_fields(self):
        _, item = self.extract(dict(self.FULL))
        self.login('kf-v')
        body = self.client.get(f'/api/knowledge/{item.id}/').json()
        for field in ('environment', 'diagnosis', 'verification', 'caveats', 'related_refs'):
            self.assertIn(field, body)

    def test_agent_search_matches_the_new_fields(self):
        """검증 방법·주의사항에만 있는 명령어가 검색에서 빠지면 지식을 못 찾는다."""
        self.extract(dict(self.FULL))
        results = json.loads(help_agent._search_knowledge(query='플래핑'))
        self.assertEqual(results['count'], 1)
        results = json.loads(help_agent._search_knowledge(query='ACOS-104904'))
        self.assertEqual(results['count'], 1)


# EVE-NG가 돌려주는 원본 모양 — 클라이언트가 이걸 우리 용어로 번역하는지 본다.
EVE_NODES = {
    '3': {'name': 'A10_1', 'template': 'a10', 'image': 'a10-vThunder-6.0.8',
          'icon': 'F5_LB.png', 'left': 150, 'top': 243, 'ram': 8192, 'cpu': 4,
          'ethernet': 3, 'url': 'telnet://eve:32771', 'status': 0},
    '1': {'name': 'Arista_1', 'template': 'veos', 'image': 'veos', 'icon': 'sw.svg',
          'left': 330, 'top': 381, 'ram': 4096, 'cpu': 1, 'ethernet': 9,
          'url': 'telnet://eve:32769', 'status': 2},
}
EVE_NETWORKS = {'5': {'name': 'Net', 'type': 'pnet0', 'left': 10, 'top': 20}}
EVE_TOPOLOGY = [
    {'type': 'ethernet', 'source': 'node3', 'source_type': 'node', 'source_label': 'E1',
     'destination': 'node1', 'destination_type': 'node', 'destination_label': 'Eth1'},
    {'type': 'ethernet', 'source': 'node1', 'source_type': 'node', 'source_label': 'Mgmt1',
     'destination': 'network5', 'destination_type': 'network', 'destination_label': ''},
]


def fake_eveng(**overrides):
    """EvengClient를 대신할 목. 실제 EVE-NG를 부르지 않는다."""
    client = MagicMock()
    client.server_version.return_value = overrides.get('version', '6.2.0-4')
    client.list_labs.return_value = overrides.get('labs', [
        {'path': '/AI-LAB-A10-OneArm.unl', 'file': 'AI-LAB-A10-OneArm.unl'},
        {'path': '/LAB_Other.unl', 'file': 'LAB_Other.unl'},
    ])
    real = eveng.EvengClient.topology
    client.topology.side_effect = lambda path: real(
        SimpleNamespace(get=lambda p: (EVE_NODES if p.endswith('/nodes')
                                       else EVE_NETWORKS if p.endswith('/networks')
                                       else EVE_TOPOLOGY)),
        path)
    return client


class EvengClientTests(TestCase):
    """EVE-NG 응답을 우리 용어로 번역하는 부분 — 원본 스키마가 밖으로 새지 않아야 한다."""

    def topology(self):
        return fake_eveng().topology('/x.unl')

    def test_translates_nodes_keyed_by_name(self):
        """이름이 키다. eve_id·console 포트는 서버를 옮기면 재부여되는 값이다."""
        nodes = {n['name']: n for n in self.topology()['nodes']}

        self.assertEqual(set(nodes), {'A10_1', 'Arista_1'})
        self.assertEqual(nodes['A10_1']['eve_id'], 3)
        self.assertEqual((nodes['A10_1']['left'], nodes['A10_1']['top']), (150, 243))
        # EVE-NG 원본 키(status, url)는 우리 이름으로 바뀌어 나간다
        self.assertNotIn('status', nodes['A10_1'])
        self.assertNotIn('url', nodes['A10_1'])

    def test_status_becomes_running_not_ready(self):
        """EVE-NG status는 프로세스가 떴다는 뜻일 뿐 부팅 완료가 아니다."""
        nodes = {n['name']: n for n in self.topology()['nodes']}
        self.assertFalse(nodes['A10_1']['running'])   # status 0
        self.assertTrue(nodes['Arista_1']['running'])  # status 2

    def test_links_use_names_and_keep_network_links(self):
        links = self.topology()['links']
        node_links = [l for l in links
                      if not l['source_is_network'] and not l['target_is_network']]

        self.assertEqual(len(links), 2)
        self.assertEqual(len(node_links), 1)
        self.assertEqual((node_links[0]['source'], node_links[0]['source_port']),
                         ('A10_1', 'E1'))
        self.assertEqual((node_links[0]['target'], node_links[0]['target_port']),
                         ('Arista_1', 'Eth1'))
        # 관리망 연결도 버리지 않는다 — 준비 판정이 관리망을 통해 이뤄진다
        net_link = [l for l in links if l['target_is_network']][0]
        self.assertEqual(net_link['target'], 'Net')

    @override_settings(EVENG_URL='', EVENG_USER='', EVENG_PASSWORD='')
    def test_missing_settings_raise_not_configured(self):
        with self.assertRaises(eveng.EvengNotConfigured):
            eveng.EvengClient()


@override_settings(EVENG_URL='http://eve.test', EVENG_USER='admin', EVENG_PASSWORD='pw')
class LabRegistryTests(TestCase):
    """랩 등록·조회·토폴로지 갱신."""

    def setUp(self):
        from .permissions import set_user_role
        for username, role in (('lr-v', 'viewer'), ('lr-e', 'engineer'), ('lr-a', 'admin')):
            user = User.objects.create_user(username, password='lr-pass-1!')
            set_user_role(user, role)

    def login(self, username):
        self.client.post('/api/auth/login/',
                         {'username': username, 'password': 'lr-pass-1!'},
                         content_type='application/json')

    def register(self, path='/AI-LAB-A10-OneArm.unl', **extra):
        return self.client.post('/api/labs/register/',
                                {'path': path, 'name': 'One-Arm', **extra},
                                content_type='application/json')

    def test_available_marks_already_registered(self):
        self.login('lr-a')
        with patch('api.views.eveng.EvengClient', return_value=fake_eveng()):
            self.register()
            body = self.client.get('/api/labs/available/').json()

        by_path = {l['path']: l for l in body['labs']}
        self.assertTrue(by_path['/AI-LAB-A10-OneArm.unl']['registered'])
        self.assertFalse(by_path['/LAB_Other.unl']['registered'])
        self.assertEqual(body['version'], '6.2.0-4')  # 서버 버전을 기록해둔다

    def test_only_registered_labs_appear_in_the_menu(self):
        """EVE-NG에는 다른 사람 작업용 랩이 섞여 있어 전부 노출하지 않는다."""
        self.login('lr-a')
        with patch('api.views.eveng.EvengClient', return_value=fake_eveng()):
            self.register()
        self.login('lr-e')
        labs = self.client.get('/api/labs/').json()

        self.assertEqual([l['path'] for l in labs], ['/AI-LAB-A10-OneArm.unl'])

    def test_refresh_stores_snapshot(self):
        self.login('lr-a')
        with patch('api.views.eveng.EvengClient', return_value=fake_eveng()):
            lab_id = self.register().json()['id']
            body = self.client.post(f'/api/labs/{lab_id}/refresh/').json()

        self.assertEqual({n['name'] for n in body['nodes']}, {'A10_1', 'Arista_1'})
        self.assertEqual(len(body['links']), 2)
        self.assertEqual(len(body['networks']), 1)
        self.assertIsNotNone(body['topology_synced_at'])

    def test_refresh_drops_nodes_removed_from_eveng(self):
        """랩에서 노드를 지우면 스냅샷에서도 사라져야 한다."""
        from .models import Lab
        self.login('lr-a')
        with patch('api.views.eveng.EvengClient', return_value=fake_eveng()):
            lab_id = self.register().json()['id']
            self.client.post(f'/api/labs/{lab_id}/refresh/')

        smaller = fake_eveng()
        smaller.topology.side_effect = None
        smaller.topology.return_value = {
            'nodes': [{'name': 'A10_1', 'eve_id': 3, 'template': 'a10', 'image': '',
                       'icon': '', 'left': 1, 'top': 2, 'ram': 0, 'cpu': 0,
                       'ethernet': 0, 'console_url': '', 'running': False}],
            'networks': [], 'links': [],
        }
        with patch('api.views.eveng.EvengClient', return_value=smaller):
            body = self.client.post(f'/api/labs/{lab_id}/refresh/').json()

        self.assertEqual([n['name'] for n in body['nodes']], ['A10_1'])
        self.assertEqual(Lab.objects.get(id=lab_id).nodes.count(), 1)

    def test_unregister_keeps_eveng_untouched(self):
        """등록 해제는 우리 기록만 지운다 — EVE-NG 삭제 API를 부르지 않는다."""
        from .models import Lab
        self.login('lr-a')
        client = fake_eveng()
        with patch('api.views.eveng.EvengClient', return_value=client):
            lab_id = self.register().json()['id']
            res = self.client.delete('/api/labs/register/', {'id': lab_id},
                                     content_type='application/json')

        self.assertEqual(res.status_code, 200)
        self.assertEqual(Lab.objects.count(), 0)
        client.delete_lab.assert_not_called()

    def test_duplicate_registration_is_rejected(self):
        self.login('lr-a')
        with patch('api.views.eveng.EvengClient', return_value=fake_eveng()):
            self.register()
            self.assertEqual(self.register().status_code, 409)

    def test_engineer_cannot_register_but_can_view(self):
        """등록은 관리자, 조회·갱신은 엔지니어."""
        self.login('lr-e')
        self.assertEqual(self.register().status_code, 403)
        self.assertEqual(self.client.get('/api/labs/').status_code, 200)

    def test_viewer_is_blocked_everywhere(self):
        self.login('lr-v')
        self.assertEqual(self.client.get('/api/labs/').status_code, 403)
        self.assertEqual(self.client.get('/api/labs/available/').status_code, 403)

    @override_settings(EVENG_URL='', EVENG_USER='', EVENG_PASSWORD='')
    def test_unconfigured_server_returns_503_not_500(self):
        """랩 서버가 없어도 앱은 살아 있어야 한다 — 화면이 안내를 띄울 수 있게."""
        from .models import LabServer
        self.login('lr-a')
        res = self.client.get('/api/labs/available/')
        self.assertEqual(res.status_code, 503)
        self.assertIn('CASEFLOW_EVENG_URL', res.json()['error'])

        # 등록도 막는다 — 안 막으면 base_url이 빈 서버 행이 생겨 랩이 붕 뜬다
        self.assertEqual(self.register().status_code, 503)
        self.assertEqual(LabServer.objects.count(), 0)


class LabProbeTests(TestCase):
    """준비 판정 — EVE-NG의 running과 프로브를 합쳐 4단계로 가른다."""

    def access(self, name, driver='a10_axapi', ip='10.0.0.1'):
        from .models import Lab, LabNodeAccess, LabServer
        server = LabServer.objects.get_or_create(base_url='http://eve.test')[0]
        lab = Lab.objects.get_or_create(server=server, path='/x.unl',
                                        defaults={'name': 'x'})[0]
        return LabNodeAccess(lab=lab, node_name=name, driver=driver, mgmt_ip=ip,
                             username='u', password='p')

    def test_running_without_access_is_unknown_not_booting(self):
        """접속 정보가 없으면 '기동 중'이 아니라 '확인 불가'다 —
        영영 안 끝나는 것처럼 보이면 뭘 해야 할지 알 수 없다."""
        states = lab_probe.node_states({'A10_1': True, 'A10_2': False}, [])
        self.assertEqual(states, {'A10_1': lab_probe.UNKNOWN, 'A10_2': lab_probe.OFF})

    def test_probe_result_splits_ready_and_booting(self):
        accesses = [self.access('A10_1'), self.access('A10_2')]
        with patch.object(lab_probe, 'probe',
                          side_effect=lambda a: a.node_name == 'A10_1'):
            states = lab_probe.node_states({'A10_1': True, 'A10_2': True}, accesses)
        self.assertEqual(states, {'A10_1': lab_probe.READY, 'A10_2': lab_probe.BOOTING})

    def test_stopped_nodes_are_not_probed(self):
        """꺼진 노드를 찌르면 타임아웃만 기다리게 된다."""
        with patch.object(lab_probe, 'probe') as mocked:
            lab_probe.node_states({'A10_1': False}, [self.access('A10_1')])
        mocked.assert_not_called()

    def test_probe_without_ip_or_driver_is_unknown(self):
        cases = [self.access('n1', driver='none'), self.access('n2', ip='')]
        for access in cases:
            self.assertIsNone(lab_probe.probe(access))

    def test_probe_failure_is_not_ready_not_crash(self):
        """장비가 부팅 중이면 예외가 난다 — 그냥 '아직 아니다'로 본다."""
        access = self.access('n1', driver='linux_ssh', ip='203.0.113.1')
        with patch('api.services.lab_probe.socket.create_connection',
                   side_effect=OSError('timed out')):
            self.assertFalse(lab_probe.probe(access))


@override_settings(EVENG_URL='http://eve.test', EVENG_USER='admin', EVENG_PASSWORD='pw')
class LabPowerAndAccessTests(TestCase):
    """전원 제어와 노드 접속 정보."""

    def setUp(self):
        from .models import Lab, LabNode, LabServer
        from .permissions import set_user_role
        for username, role in (('lp-v', 'viewer'), ('lp-e', 'engineer')):
            user = User.objects.create_user(username, password='lp-pass-1!')
            set_user_role(user, role)
        server = LabServer.objects.create(base_url='http://eve.test')
        self.lab = Lab.objects.create(server=server, path='/x.unl', name='x')
        LabNode.objects.create(lab=self.lab, name='A10_1', eve_id=3, ram=8192)
        LabNode.objects.create(lab=self.lab, name='Arista_1', eve_id=1, ram=4096)

    def login(self, username='lp-e'):
        self.client.post('/api/auth/login/',
                         {'username': username, 'password': 'lp-pass-1!'},
                         content_type='application/json')

    def test_power_starts_heaviest_first(self):
        """한꺼번에 켜면 공용 EVE-NG가 부담을 받는다 — 무거운 것부터 순차로."""
        client = MagicMock()
        self.login()
        with patch('api.views.eveng.EvengClient', return_value=client), \
             patch('api.views.time.sleep'), \
             patch('api.views.connection.close'), \
             patch('api.views.threading.Thread',
                   side_effect=lambda target, args, daemon: SimpleNamespace(
                       start=lambda: target(*args))):
            res = self.client.post(f'/api/labs/{self.lab.id}/power/', {'action': 'start'},
                                   content_type='application/json')

        self.assertEqual(res.status_code, 202)
        started = [call.args[1] for call in client.start_node.call_args_list]
        self.assertEqual(started, [3, 1])  # A10(8GB) -> Arista(4GB)

    def test_power_rejects_unknown_action(self):
        self.login()
        res = self.client.post(f'/api/labs/{self.lab.id}/power/', {'action': 'reboot'},
                               content_type='application/json')
        self.assertEqual(res.status_code, 400)

    def test_power_ignores_nodes_not_in_the_lab(self):
        """오타로 남의 랩 노드를 끄는 일이 없게 토폴로지에 있는 이름만 받는다."""
        self.login()
        with patch('api.views._power_worker'):
            res = self.client.post(f'/api/labs/{self.lab.id}/power/',
                                   {'action': 'stop', 'nodes': ['A10_1', '남의노드']},
                                   content_type='application/json')
        self.assertEqual(res.json()['nodes'], ['A10_1'])

    def test_viewer_cannot_control_power(self):
        self.login('lp-v')
        res = self.client.post(f'/api/labs/{self.lab.id}/power/', {'action': 'stop'},
                               content_type='application/json')
        self.assertEqual(res.status_code, 403)

    def test_access_password_is_never_returned(self):
        self.login()
        self.client.put(f'/api/labs/{self.lab.id}/access/', {'rows': [
            {'node_name': 'A10_1', 'mgmt_ip': '10.0.0.1', 'driver': 'a10_axapi',
             'username': 'admin', 'password': 'lab-secret'},
        ]}, content_type='application/json')

        body = self.client.get(f'/api/labs/{self.lab.id}/access/').json()
        row = [r for r in body if r['node_name'] == 'A10_1'][0]
        self.assertTrue(row['has_password'])
        self.assertNotIn('lab-secret', json.dumps(body))
        self.assertNotIn('password', row)

    def test_blank_password_keeps_the_stored_one(self):
        """화면은 비밀번호를 안 받아오므로 매번 빈 값으로 올라온다 —
        그걸로 덮으면 저장해둔 값이 날아간다."""
        from .models import LabNodeAccess
        self.login()
        rows = [{'node_name': 'A10_1', 'mgmt_ip': '10.0.0.1', 'driver': 'a10_axapi',
                 'username': 'admin', 'password': 'keep-me'}]
        self.client.put(f'/api/labs/{self.lab.id}/access/', {'rows': rows},
                        content_type='application/json')
        rows[0]['password'] = ''
        rows[0]['role'] = 'lb-primary'
        self.client.put(f'/api/labs/{self.lab.id}/access/', {'rows': rows},
                        content_type='application/json')

        access = LabNodeAccess.objects.get(lab=self.lab, node_name='A10_1')
        self.assertEqual(access.password, 'keep-me')
        self.assertEqual(access.role, 'lb-primary')

    def test_access_rejects_names_not_in_topology(self):
        """오타로 만든 유령 행이 남으면 나중에 어느 노드 얘기인지 알 수 없다."""
        from .models import LabNodeAccess
        self.login()
        self.client.put(f'/api/labs/{self.lab.id}/access/', {'rows': [
            {'node_name': '없는노드', 'mgmt_ip': '10.0.0.9', 'driver': 'a10_axapi'},
        ]}, content_type='application/json')
        self.assertEqual(LabNodeAccess.objects.count(), 0)

    def test_access_survives_topology_refresh(self):
        """사람이 적은 값은 EVE-NG 갱신으로 덮이면 안 된다."""
        from .models import LabNodeAccess
        self.login()
        self.client.put(f'/api/labs/{self.lab.id}/access/', {'rows': [
            {'node_name': 'A10_1', 'mgmt_ip': '10.0.0.1', 'driver': 'a10_axapi',
             'username': 'admin', 'password': 'pw'},
        ]}, content_type='application/json')

        with patch('api.views.eveng.EvengClient', return_value=fake_eveng()):
            self.client.post(f'/api/labs/{self.lab.id}/refresh/')

        access = LabNodeAccess.objects.get(lab=self.lab, node_name='A10_1')
        self.assertEqual((access.mgmt_ip, access.password), ('10.0.0.1', 'pw'))

    def test_status_counts_and_lists_unprobeable(self):
        self.login()
        with patch('api.views.eveng.EvengClient') as client:
            client.return_value.node_states.return_value = {'A10_1': True, 'Arista_1': False}
            body = self.client.get(f'/api/labs/{self.lab.id}/status/').json()

        self.assertEqual(body['counts'], {'unknown': 1, 'off': 1})
        self.assertEqual(body['unprobeable'], ['A10_1'])
        self.assertEqual(body['total'], 2)


class PortNormalizeTests(TestCase):
    """같은 포트를 EVE-NG는 'E1', 장비는 'Ethernet1', LLDP는 'Et1'이라 부른다."""

    def test_same_port_written_differently_matches(self):
        forms = ['E1', 'Eth1', 'Ethernet1', 'et1', ' Ethernet 1 ']
        self.assertEqual({lab_drivers.normalize_port(f) for f in forms}, {'e1'})

    def test_different_ports_stay_different(self):
        self.assertNotEqual(lab_drivers.normalize_port('Eth1'),
                            lab_drivers.normalize_port('Eth2'))
        self.assertNotEqual(lab_drivers.normalize_port('Eth1'),
                            lab_drivers.normalize_port('Mgmt1'))

    def test_unparseable_name_is_kept_as_is(self):
        self.assertEqual(lab_drivers.normalize_port('port-a/b'), 'port-a/b')


class LabDriverTests(TestCase):
    """드라이버는 벤더 응답을 공통 모양으로 바꾼다."""

    def access(self, driver='arista_eapi'):
        return SimpleNamespace(driver=driver, mgmt_ip='10.0.0.1',
                               username='admin', password='pw', node_name='Arista_1')

    def test_arista_parses_hostname_and_lldp(self):
        payload = {'result': [{'lldpNeighbors': [
            {'port': 'Ethernet1', 'neighborDevice': 'A10_1.lab', 'neighborPort': 'E1'},
        ]}]}
        driver = lab_drivers.AristaDriver(self.access())
        with patch('api.services.lab_drivers.requests.post') as post:
            post.return_value = MagicMock(json=lambda: payload)
            neighbors = driver.lldp_neighbors()
        self.assertEqual(neighbors, [{'local_port': 'Ethernet1',
                                      'remote_host': 'A10_1.lab', 'remote_port': 'E1'}])

    def test_arista_error_in_200_body_is_raised(self):
        """eAPI는 200으로 오류를 돌려준다 — 상태 코드만 보면 놓친다."""
        driver = lab_drivers.AristaDriver(self.access())
        with patch('api.services.lab_drivers.requests.post') as post:
            post.return_value = MagicMock(
                json=lambda: {'error': {'message': 'invalid command'}})
            with self.assertRaises(lab_drivers.DriverError):
                driver.hostname()

    def test_a10_session_always_sends_json_content_type(self):
        """DELETE에도 이 헤더가 없으면 415다(실측). 세션에 박아둔다."""
        driver = lab_drivers.A10Driver(self.access('a10_axapi'))
        self.assertEqual(driver.session.headers['Content-Type'], 'application/json')

    def test_a10_auth_failure_is_reported(self):
        driver = lab_drivers.A10Driver(self.access('a10_axapi'))
        with patch.object(driver.session, 'post') as post:
            post.return_value = MagicMock(status_code=401, json=lambda: {})
            with self.assertRaises(lab_drivers.DriverError):
                driver.hostname()

    def test_output_is_truncated(self):
        """show tech는 수 MB다. 판정은 코드가 하니 원문 전체가 필요 없다."""
        text = lab_drivers.truncate('x' * (lab_drivers.MAX_OUTPUT_CHARS + 500))
        self.assertIn('이하 생략', text)
        self.assertLess(len(text), lab_drivers.MAX_OUTPUT_CHARS + 200)

    def test_driver_is_none_without_ip_or_support(self):
        self.assertIsNone(lab_drivers.get_driver(
            SimpleNamespace(driver='arista_eapi', mgmt_ip='')))
        self.assertIsNone(lab_drivers.get_driver(
            SimpleNamespace(driver='none', mgmt_ip='10.0.0.1')))


@override_settings(EVENG_URL='http://eve.test', EVENG_USER='admin', EVENG_PASSWORD='pw')
class LabCheckTests(TestCase):
    """읽기 전용 점검 — 어긋난 것을 실제로 잡아내는지가 전부다."""

    def setUp(self):
        from .models import Lab, LabLink, LabNode, LabNodeAccess, LabServer
        from .permissions import set_user_role
        user = User.objects.create_user('lc-e', password='lc-pass-1!')
        set_user_role(user, 'engineer')
        server = LabServer.objects.create(base_url='http://eve.test')
        self.lab = Lab.objects.create(server=server, path='/x.unl', name='x')
        for name in ('Arista_1', 'A10_1'):
            LabNode.objects.create(lab=self.lab, name=name, eve_id=1)
            LabNodeAccess.objects.create(lab=self.lab, node_name=name,
                                         mgmt_ip='10.0.0.1', driver='arista_eapi',
                                         username='u', password='p')
        # EVE-NG가 아는 배선: Arista_1:Eth1 <-> A10_1:E1
        LabLink.objects.create(lab=self.lab, source='Arista_1', source_port='Eth1',
                               target='A10_1', target_port='E1')
        # 관리망 연결은 LLDP 대조 대상이 아니다
        LabLink.objects.create(lab=self.lab, source='Arista_1', source_port='Mgmt1',
                               target='Net', target_port='', target_is_network=True)

    def login(self):
        self.client.post('/api/auth/login/',
                         {'username': 'lc-e', 'password': 'lc-pass-1!'},
                         content_type='application/json')

    def run_with(self, facts_by_node):
        def fake_facts(access):
            return {'access': access, **facts_by_node[access.node_name]}
        with patch('api.services.lab_check._facts', side_effect=fake_facts):
            self.login()
            return self.client.post(f'/api/labs/{self.lab.id}/check/').json()

    def matching_facts(self):
        return {
            'Arista_1': {'hostname': 'Arista_1', 'neighbors': [
                {'local_port': 'Ethernet1', 'remote_host': 'A10_1', 'remote_port': 'E1'}]},
            'A10_1': {'hostname': 'A10_1', 'neighbors': [
                {'local_port': 'E1', 'remote_host': 'Arista_1', 'remote_port': 'Ethernet1'}]},
        }

    def test_matching_wiring_passes(self):
        body = self.run_with(self.matching_facts())
        self.assertEqual(body['counts']['fail'], 0)
        wiring = [r for r in body['results'] if r['check'] == '배선 대조']
        self.assertTrue(all(r['status'] == 'pass' for r in wiring), wiring)

    def test_mismatched_wiring_is_caught(self):
        """이 단계의 완료 기준 — 배선을 일부러 어긋나게 하면 실패로 잡아야 한다.
        통과만 되는 검증은 검증이 아니다."""
        facts = self.matching_facts()
        # Arista_1이 Eth1이 아니라 Eth9에서 이웃을 본다고 하자
        facts['Arista_1']['neighbors'] = [
            {'local_port': 'Ethernet9', 'remote_host': 'A10_1', 'remote_port': 'E1'}]
        body = self.run_with(facts)

        failed = [r for r in body['results']
                  if r['check'] == '배선 대조' and r['status'] == 'fail']
        self.assertEqual(len(failed), 1)
        self.assertIn('EVE-NG에는 있는데 장비가 못 봄', failed[0]['detail'])
        self.assertIn('장비는 보는데 EVE-NG에 없음', failed[0]['detail'])

    def test_wrong_device_behind_the_ip_is_caught(self):
        """관리 IP가 다른 장비를 가리키면 엉뚱한 곳에 설정이 들어간다."""
        facts = self.matching_facts()
        facts['Arista_1']['hostname'] = '남의장비'
        body = self.run_with(facts)

        failed = [r for r in body['results']
                  if r['check'] == '장비 확인' and r['status'] == 'fail']
        self.assertEqual(len(failed), 1)
        self.assertIn('남의장비', failed[0]['detail'])

    def test_port_naming_differences_do_not_cause_false_failures(self):
        """장비는 Ethernet1, EVE-NG는 Eth1이라고 적는다 — 이건 불일치가 아니다."""
        body = self.run_with(self.matching_facts())
        self.assertEqual(body['counts']['fail'], 0)

    def test_unreadable_lldp_is_skip_not_fail(self):
        """LLDP를 못 읽는 것과 배선이 틀린 것은 다르다."""
        facts = self.matching_facts()
        facts['A10_1']['neighbors'] = None
        body = self.run_with(facts)

        wiring = {r['node']: r['status'] for r in body['results']
                  if r['check'] == '배선 대조'}
        self.assertEqual(wiring['A10_1'], 'skip')
        self.assertEqual(wiring['Arista_1'], 'pass')

    def test_nodes_without_access_are_reported_as_skip(self):
        from .models import LabNode
        LabNode.objects.create(lab=self.lab, name='Server_1', eve_id=6)
        body = self.run_with(self.matching_facts())

        skipped = [r for r in body['results']
                   if r['check'] == '접속 정보' and r['node'] == 'Server_1']
        self.assertEqual(len(skipped), 1)
        self.assertEqual(skipped[0]['status'], 'skip')


class LabConfigTests(TestCase):
    """Lab Tests — EVE-NG 설정 여부 조회. 랩 서버가 없어도 앱은 정상 동작해야 한다."""

    def setUp(self):
        from .permissions import set_user_role
        for username, role in (('lab-v', 'viewer'), ('lab-e', 'engineer')):
            user = User.objects.create_user(username, password='lab-pass-1!')
            set_user_role(user, role)

    def login(self, username):
        self.client.post('/api/auth/login/',
                         {'username': username, 'password': 'lab-pass-1!'},
                         content_type='application/json')

    @override_settings(EVENG_URL='', EVENG_USER='', EVENG_PASSWORD='')
    def test_reports_unconfigured_without_failing(self):
        """설정이 없어도 500이 아니라 configured=false를 돌려줘야 화면이 안내를 띄운다."""
        self.login('lab-e')
        res = self.client.get('/api/labs/config/')
        self.assertEqual(res.status_code, 200)
        self.assertEqual(res.json(), {'configured': False, 'server': ''})

    @override_settings(EVENG_URL='http://10.0.0.5', EVENG_USER='admin',
                       EVENG_PASSWORD='secret')
    def test_reports_configured_without_leaking_credentials(self):
        self.login('lab-e')
        body = self.client.get('/api/labs/config/').json()

        self.assertTrue(body['configured'])
        self.assertEqual(body['server'], 'http://10.0.0.5')
        # 계정·비밀번호는 어떤 형태로도 나가지 않는다
        self.assertNotIn('admin', json.dumps(body))
        self.assertNotIn('secret', json.dumps(body))

    @override_settings(EVENG_URL='http://10.0.0.5', EVENG_USER='admin', EVENG_PASSWORD='')
    def test_partial_settings_count_as_unconfigured(self):
        """셋 중 하나라도 비면 접속이 안 되므로 설정된 것으로 보지 않는다."""
        self.login('lab-e')
        self.assertFalse(self.client.get('/api/labs/config/').json()['configured'])

    def test_viewer_is_blocked(self):
        """노드 전원 제어가 엔지니어 이상이라 랩 화면 자체를 같은 기준으로 막는다."""
        self.login('lab-v')
        self.assertEqual(self.client.get('/api/labs/config/').status_code, 403)


class KnowledgeModelSettingTests(TestCase):
    """지식 추출 모델은 메일 분석 모델과 분리돼 있고, 상위 두 모델로 제한된다."""

    def setUp(self):
        from .permissions import set_user_role
        for username, role in (('km-v', 'viewer'), ('km-e', 'engineer'), ('km-a', 'admin')):
            user = User.objects.create_user(username, password='km-pass-1!')
            set_user_role(user, role)

    def login(self, username):
        self.client.post('/api/auth/login/',
                         {'username': username, 'password': 'km-pass-1!'},
                         content_type='application/json')

    def test_defaults_to_opus5_and_ignores_unknown_stored_value(self):
        self.assertEqual(analyzer.get_knowledge_model(), 'claude-opus-5')
        # 카탈로그에서 빠진 모델이 저장돼 있어도 지식 추출이 죽으면 안 된다
        AppSetting.set(analyzer.KNOWLEDGE_MODEL_SETTING_KEY, 'gemini-3.5-flash')
        self.assertEqual(analyzer.get_knowledge_model(), 'claude-opus-5')

    def test_candidates_never_fall_back_to_cheap_models(self):
        """저비용 모델로 조용히 떨어지면 눈에 안 띄는 품질 저하가 지식에 남는다."""
        AppSetting.set(analyzer.KNOWLEDGE_MODEL_SETTING_KEY, 'claude-sonnet-5')
        self.assertEqual(analyzer.knowledge_model_candidates(),
                         ['claude-sonnet-5', 'claude-opus-5'])
        self.assertNotIn('claude-haiku-4-5', analyzer.knowledge_model_candidates())

    def test_is_independent_of_translation_model(self):
        """메일 분석 모델을 바꿔도 지식 추출 모델은 그대로여야 한다."""
        AppSetting.set(analyzer.TRANSLATION_MODEL_SETTING_KEY, 'claude-haiku-4-5')
        self.assertEqual(analyzer.get_translation_model(), 'claude-haiku-4-5')
        self.assertEqual(analyzer.get_knowledge_model(), 'claude-opus-5')

    def test_get_is_open_and_put_is_admin_only(self):
        self.login('km-v')
        res = self.client.get('/api/settings/knowledge-model/')
        self.assertEqual(res.status_code, 200)
        self.assertEqual(res.json()['current'], 'claude-opus-5')
        self.assertEqual([m['id'] for m in res.json()['models']],
                         ['claude-opus-5', 'claude-sonnet-5'])

        self.login('km-e')
        self.assertEqual(
            self.client.put('/api/settings/knowledge-model/', {'model': 'claude-sonnet-5'},
                            content_type='application/json').status_code, 403)

    def test_analyzed_by_records_the_model_that_actually_answered(self):
        """폴백이 걸리면 설정 모델과 답한 모델이 갈린다 — 기록은 답한 쪽이어야 한다."""
        from .services import knowledge
        case = make_case(vendor='A10', status='Resolved', resolution='조치')
        result = {'has_knowledge': True, 'title': '제목', 'problem': '문제',
                  'root_cause': '원인', 'resolution': '조치', 'device_model': '',
                  'software_version': ''}
        with patch.object(knowledge, 'generate_structured_with_model',
                          return_value=('claude-sonnet-5', result)), \
             patch.object(knowledge, 'enrich_with_references'):
            outcome, item = knowledge.extract_knowledge(case)

        self.assertEqual(outcome, 'created')
        self.assertEqual(analyzer.get_knowledge_model(), 'claude-opus-5')  # 설정값
        self.assertEqual(item.analyzed_by, 'claude-sonnet-5')             # 실제 답한 모델

    @override_settings(ANTHROPIC_API_KEY='test-key')
    def test_admin_can_switch_between_the_two_models_only(self):
        self.login('km-a')
        res = self.client.put('/api/settings/knowledge-model/', {'model': 'claude-sonnet-5'},
                              content_type='application/json')
        self.assertEqual(res.status_code, 200)
        self.assertEqual(res.json()['current'], 'claude-sonnet-5')
        self.assertEqual(analyzer.get_knowledge_model(), 'claude-sonnet-5')

        # 목록 밖 모델은 거부 — 지식 품질이 조용히 떨어지는 경로를 막는다
        res = self.client.put('/api/settings/knowledge-model/', {'model': 'claude-haiku-4-5'},
                              content_type='application/json')
        self.assertEqual(res.status_code, 400)
        self.assertEqual(analyzer.get_knowledge_model(), 'claude-sonnet-5')


class KnowledgeBaseTests(TestCase):
    """지식 베이스 — 추출 서비스 필터링, API 권한, 에이전트 검색 도구."""

    def setUp(self):
        from .permissions import set_user_role
        for username, role in (('kv1', 'viewer'), ('ke1', 'engineer'), ('ka1', 'admin')):
            user = User.objects.create_user(username, password='knowledge-pass-1!')
            set_user_role(user, role)
        self.case = make_case(vendor='A10', status='Resolved',
                              summary='SSL RST 케이스', resolution='P14 업그레이드')
        make_email(self.case, 'Re: Case # 1 SSL RST')

    def login(self, username):
        self.client.post('/api/auth/login/',
                         {'username': username, 'password': 'knowledge-pass-1!'},
                         content_type='application/json')

    def make_item(self, **kwargs):
        from .models import KnowledgeItem
        defaults = dict(case=self.case, vendor='A10', title='SSL RST 해결',
                        problem='RST 발생', resolution='ACOS 5.2.1-P14 업그레이드')
        defaults.update(kwargs)
        return KnowledgeItem.objects.create(**defaults)

    def test_extract_saves_draft_with_case_fallback_fields(self):
        from .services import knowledge
        self.case.device_model = 'TH5440S'
        self.case.save()
        result = {'has_knowledge': True, 'title': '제목', 'problem': '문제',
                  'root_cause': '원인', 'resolution': 'CLI 조치',
                  'device_model': '', 'software_version': '5.2.1-P7'}
        with patch.object(knowledge, 'generate_structured_with_model',
                          return_value=(MODEL, result)):
            outcome, item = knowledge.extract_knowledge(self.case)
        self.assertEqual(outcome, 'created')
        self.assertEqual(item.status, 'draft')
        # AI가 빈 값을 준 필드는 케이스 값으로 폴백
        self.assertEqual(item.device_model, 'TH5440S')
        self.assertEqual(item.software_version, '5.2.1-P7')

    def test_extract_skips_no_knowledge_and_existing(self):
        from .services import knowledge
        no_knowledge = {'has_knowledge': False, 'title': '', 'problem': '',
                        'root_cause': '', 'resolution': '',
                        'device_model': '', 'software_version': ''}
        with patch.object(knowledge, 'generate_structured_with_model',
                          return_value=(MODEL, no_knowledge)):
            outcome, item = knowledge.extract_knowledge(self.case)
        self.assertEqual((outcome, item), ('no_knowledge', None))

        existing = self.make_item()
        with patch.object(knowledge, 'generate_structured_with_model') as mocked:
            outcome, item = knowledge.extract_knowledge(self.case)
        mocked.assert_not_called()  # 기존 항목이 있으면 AI 호출 자체를 안 함
        self.assertEqual((outcome, item), ('exists', existing))

    def test_api_roles_and_confirm_flow(self):
        item = self.make_item()
        url = f'/api/knowledge/{item.id}/'

        self.login('kv1')  # viewer: 조회만
        self.assertEqual(self.client.get('/api/knowledge/').status_code, 200)
        self.assertEqual(self.client.patch(url, {'status': 'confirmed'},
                                           content_type='application/json').status_code, 403)
        self.assertEqual(self.client.delete(url).status_code, 403)

        self.login('ke1')  # engineer: 수정·확정 가능, 삭제 불가
        res = self.client.patch(url, {'status': 'confirmed', 'title': '수정된 제목'},
                                content_type='application/json')
        self.assertEqual(res.status_code, 200)
        self.assertEqual(res.json()['status'], 'confirmed')
        self.assertEqual(res.json()['title'], '수정된 제목')
        self.assertEqual(res.json()['source_case']['case_id'], self.case.case_id)
        self.assertEqual(self.client.delete(url).status_code, 403)

        self.login('ka1')  # admin: 삭제 가능
        self.assertEqual(self.client.delete(url).status_code, 204)

    def test_enrich_validates_indexes_and_saves_references(self):
        from .services import knowledge
        from .services import references as refdocs
        item = self.make_item()
        candidates = [
            {'document': 'A10/guide.pdf', 'pages': 'p.1-2', 'score': 0.9,
             'text': 'ssl 설정 섹션', 'title': '', 'vendor': 'A10'},
            {'document': 'A10/guide.pdf', 'pages': 'p.9-10', 'score': 0.5,
             'text': '무관 섹션', 'title': '', 'vendor': 'A10'},
        ]
        # AI가 유효 index 0과 존재하지 않는 index 7을 반환 → 7은 코드 검증에서 버려짐
        ai_result = {'relevant': [{'index': 0, 'note': '해결 절차 근거'},
                                  {'index': 7, 'note': '지어낸 인용'}]}
        with patch.object(refdocs, 'search', return_value=candidates), \
             patch.object(knowledge, 'generate_structured', return_value=ai_result):
            outcome = knowledge.enrich_with_references(item)
        self.assertEqual(outcome, 'enriched')
        item.refresh_from_db()
        self.assertEqual(len(item.references), 1)
        self.assertEqual(item.references[0]['pages'], 'p.1-2')
        self.assertEqual(item.references[0]['note'], '해결 절차 근거')

        # 후보가 없으면 references는 빈 목록으로 확정
        with patch.object(refdocs, 'search', return_value=[]):
            self.assertEqual(knowledge.enrich_with_references(item), 'no_candidates')
        item.refresh_from_db()
        self.assertEqual(item.references, [])

    def test_search_knowledge_tool_filters_and_prefers_confirmed(self):
        self.make_item(title='VRRP 페일오버 반복', resolution='preempt 설정 수정',
                       status='confirmed')
        self.make_item(title='VRRP 로그 문의', resolution='로그 레벨 조정')
        result = json.loads(help_agent._search_knowledge('VRRP'))
        self.assertEqual(result['count'], 2)
        self.assertEqual(result['results'][0]['status'], 'confirmed')  # 확정 우선
        self.assertEqual(result['results'][0]['source_case'], self.case.case_id)
        # 본문(해결 조치) 키워드로도 검색된다
        self.assertEqual(json.loads(help_agent._search_knowledge('preempt'))['count'], 1)
        self.assertEqual(json.loads(help_agent._search_knowledge('없는키워드'))['count'], 0)


class ReferenceSearchTests(TestCase):
    """레퍼런스 문서 — 청킹, 해시 캐싱 인제스트, 벡터 검색 (임베딩 API는 모킹)."""

    def setUp(self):
        from .services import references
        references._invalidate_cache()

    def make_doc_with_chunks(self, vendor, filename, vectors_and_texts):
        import numpy as np
        from django.conf import settings
        from .models import ReferenceChunk, ReferenceDocument
        doc = ReferenceDocument.objects.create(
            vendor=vendor, filename=filename, sha256='x' * 64,
            embedding_model=settings.EMBEDDING_MODEL, chunk_count=len(vectors_and_texts))
        for seq, (vec, text) in enumerate(vectors_and_texts):
            ReferenceChunk.objects.create(
                document=doc, seq=seq, page_start=seq + 1, page_end=seq + 1,
                text=text, embedding=np.asarray(vec, dtype=np.float32).tobytes(),
                embedding_model=settings.EMBEDDING_MODEL)
        return doc

    def test_chunk_pages_tracks_page_ranges_and_overlap(self):
        from .services import references
        pages = [(1, 'a' * 3000), (2, 'b' * 3000), (3, 'c' * 500)]
        chunks = references.chunk_pages(pages)
        self.assertGreaterEqual(len(chunks), 2)
        self.assertEqual(chunks[0]['page_start'], 1)
        # 첫 청크는 1페이지를 넘어 2페이지 내용까지 포함
        self.assertEqual(chunks[0]['page_end'], 2)
        self.assertTrue(all(len(c['text']) <= references.CHUNK_CHARS for c in chunks))
        # 오버랩: 다음 청크 머리에 이전 청크 꼬리가 겹침
        self.assertTrue(chunks[1]['text'].startswith(
            chunks[0]['text'][-references.OVERLAP_CHARS + 100:][:100]))

    def test_search_ranks_by_similarity_and_filters_vendor(self):
        import numpy as np
        from .services import references
        self.make_doc_with_chunks('A10', 'A10/guide.pdf', [
            ([1.0, 0.0, 0.0, 0.0], 'SSL 오프로드 섹션'),
            ([0.0, 1.0, 0.0, 0.0], 'VRRP 섹션'),
        ])
        self.make_doc_with_chunks('Arista', 'Arista/eos.pdf', [
            ([0.9, 0.1, 0.0, 0.0], 'EOS SSL 유사 섹션'),
        ])
        query_vec = np.asarray([[1.0, 0.0, 0.0, 0.0]], dtype=np.float32)
        with patch.object(references, 'embed_texts', return_value=query_vec):
            results = references.search('ssl offload')
            self.assertEqual(results[0]['text'], 'SSL 오프로드 섹션')
            self.assertEqual(results[0]['document'], 'A10/guide.pdf')
            self.assertIn('p.1-1', results[0]['pages'])
            # 벤더 필터
            arista_only = references.search('ssl offload', vendor='Arista')
            self.assertEqual([r['vendor'] for r in arista_only], ['Arista'])

    def test_search_empty_without_ingested_docs(self):
        from .services import references
        self.assertEqual(references.search('anything'), [])

    def test_ingest_skips_unchanged_and_reprocesses_on_force(self):
        import tempfile
        from pathlib import Path
        import numpy as np
        from .models import ReferenceDocument
        from .services import references

        with tempfile.TemporaryDirectory() as tmp:
            pdf = Path(tmp) / 'guide.pdf'
            pdf.write_bytes(b'%PDF-fake')
            fake_pages = [(1, 'ACOS 6.0.8 Test Guide © A10'), (2, 'slb 설정 본문')]
            with patch.object(references, 'extract_pages', return_value=fake_pages), \
                 patch.object(references, 'embed_texts',
                              return_value=np.ones((1, 4), dtype=np.float32)) as embed:
                self.assertEqual(
                    references.ingest_file('A10', 'config', 'A10/config/guide.pdf', pdf),
                    'created')
                doc = ReferenceDocument.objects.get(filename='A10/config/guide.pdf')
                self.assertEqual(doc.title, 'ACOS 6.0.8 Test Guide')
                self.assertEqual(doc.doc_type, 'config')
                self.assertEqual(doc.chunk_count, 1)
                # 같은 파일 재실행 → 임베딩 호출 없이 건너뜀
                embed.reset_mock()
                self.assertEqual(
                    references.ingest_file('A10', 'config', 'A10/config/guide.pdf', pdf),
                    'skipped')
                embed.assert_not_called()
                # --force → 재처리
                self.assertEqual(
                    references.ingest_file('A10', 'config', 'A10/config/guide.pdf', pdf,
                                           force=True),
                    'updated')

    def test_xlsx_rows_become_chunks(self):
        import tempfile
        from pathlib import Path
        from openpyxl import Workbook
        from .services import references

        with tempfile.TemporaryDirectory() as tmp:
            xlsx = Path(tmp) / 'issues.xlsx'
            wb = Workbook()
            ws = wb.active
            ws.title = '이슈'
            ws.append(['이슈번호', '장비', '증상', '조치'])
            ws.append(['I-001', 'TH3350', 'RST 발생', 'P14 업그레이드'])
            ws.append([None, None, None, None])  # 빈 행은 무시
            ws.append(['I-002', 'AP730', None, 'RMA 진행'])  # 빈 셀은 생략
            wb.save(xlsx)

            title, chunks = references.extract_xlsx_rows(xlsx)
            self.assertEqual(title, 'issues')
            self.assertEqual(len(chunks), 2)
            self.assertIn('이슈번호: I-001', chunks[0]['text'])
            self.assertIn('[이슈 시트 2행]', chunks[0]['text'])
            self.assertIn('조치: RMA 진행', chunks[1]['text'])
            self.assertNotIn('증상:', chunks[1]['text'])  # 빈 셀 생략
            self.assertEqual(chunks[1]['page_start'], 4)  # 실제 행 번호 유지


@override_settings(TRANSLATION_FALLBACK_MODELS=['claude-haiku-4-5'],
                   ANTHROPIC_API_KEY='k', GOOGLE_API_KEY='k')
class AnalyzerFallbackTests(TestCase):
    """무료 모델이 실패하면 폴백 모델로 재시도 (429/503 대비)."""

    ANALYSIS = {'subject_ko': '제목', 'body_ko': '본문', 'summary': '요약',
                'description': '설명', 'action_update': '조치', 'resolution': '',
                'suggested_status': 'Open', 'device_model': '', 'device_serial': '',
                'software_version': ''}

    def setUp(self):
        analyzer._cooldown_until.clear()   # 할당량 쿨다운은 프로세스 전역 상태

    def analyze(self):
        return analyzer.analyze_email('subject', 'body', 'inbound', True)

    def test_fallback_model_is_used_when_free_model_fails(self):
        calls = []

        def call(model, *args, **kwargs):
            calls.append(model)
            if model.startswith('gemini'):
                raise RuntimeError('429 quota exceeded')
            return dict(self.ANALYSIS)

        with patch.object(analyzer, '_call_provider', side_effect=call), \
                analyzer.translation_model_override('gemini-3.5-flash'):
            result = self.analyze()

        self.assertEqual(calls, ['gemini-3.5-flash', 'claude-haiku-4-5'])
        self.assertEqual(result[analyzer.ANALYZED_BY_KEY], 'claude-haiku-4-5')

    def test_fallback_is_not_called_when_free_model_succeeds(self):
        with patch.object(analyzer, '_call_provider',
                          return_value=dict(self.ANALYSIS)) as mocked, \
                analyzer.translation_model_override('gemini-3.5-flash'):
            result = self.analyze()

        self.assertEqual(mocked.call_count, 1)
        self.assertEqual(result[analyzer.ANALYZED_BY_KEY], 'gemini-3.5-flash')

    def test_unparsable_response_also_falls_back(self):
        # 예외 없이 None(파싱 실패)이 와도 다음 모델로 넘어간다
        results = [None, dict(self.ANALYSIS)]
        with patch.object(analyzer, '_call_provider', side_effect=results), \
                analyzer.translation_model_override('gemini-3.5-flash'):
            result = self.analyze()

        self.assertEqual(result[analyzer.ANALYZED_BY_KEY], 'claude-haiku-4-5')

    def test_all_models_failing_returns_none(self):
        with patch.object(analyzer, '_call_provider', side_effect=RuntimeError('boom')), \
                analyzer.translation_model_override('gemini-3.5-flash'):
            self.assertIsNone(self.analyze())

    def test_case_records_the_model_that_actually_ran(self):
        case = make_case(vendor='A10')
        analysis = dict(self.ANALYSIS, **{analyzer.ANALYZED_BY_KEY: 'claude-haiku-4-5'})
        with analyzer.translation_model_override('gemini-3.5-flash'):
            gmail_sync.apply_analysis_to_case(case, analysis, 'inbound', timezone.now())
        case.refresh_from_db()
        self.assertEqual(case.analyzed_by, 'claude-haiku-4-5')

    def test_quota_error_pauses_the_model_for_later_calls(self):
        # 무료 한도가 소진되면 메일마다 다시 찔러보지 않는다
        calls = []

        def call(model, *args, **kwargs):
            calls.append(model)
            if model.startswith('gemini'):
                raise RuntimeError('429 RESOURCE_EXHAUSTED: quota exceeded')
            return dict(self.ANALYSIS)

        with patch.object(analyzer, '_call_provider', side_effect=call), \
                analyzer.translation_model_override('gemini-3.5-flash'):
            self.analyze()
            self.analyze()

        self.assertEqual(calls.count('gemini-3.5-flash'), 1)
        self.assertEqual(calls.count('claude-haiku-4-5'), 2)

    def test_non_quota_error_does_not_pause_the_model(self):
        calls = []

        def call(model, *args, **kwargs):
            calls.append(model)
            if model.startswith('gemini'):
                raise RuntimeError('connection reset')
            return dict(self.ANALYSIS)

        with patch.object(analyzer, '_call_provider', side_effect=call), \
                analyzer.translation_model_override('gemini-3.5-flash'):
            self.analyze()
            self.analyze()

        self.assertEqual(calls.count('gemini-3.5-flash'), 2)

    def test_model_override_restores_the_previous_model(self):
        with analyzer.translation_model_override('gemini-3.5-flash'):
            with analyzer.translation_model_override('claude-haiku-4-5'):
                self.assertEqual(analyzer.get_translation_model(), 'claude-haiku-4-5')
            self.assertEqual(analyzer.get_translation_model(), 'gemini-3.5-flash')


class SyncGmailCommandTests(TestCase):
    """sync_gmail 관리 명령 — cron 실행용 옵션."""

    def test_model_option_overrides_only_this_run(self):
        seen = {}

        def fake_sync(max_results):
            seen['model'] = analyzer.get_translation_model()
            seen['max_results'] = max_results
            return dict(fetched=0, cases_created=0, emails_added=0,
                        ignored=0, no_vendor=0, skipped=0, errors=0)

        with patch('api.management.commands.sync_gmail.sync_gmail', side_effect=fake_sync):
            call_command('sync_gmail', '--model', 'gemini-3.5-flash',
                         '--max-results', '20', stdout=StringIO())

        self.assertEqual(seen, {'model': 'gemini-3.5-flash', 'max_results': 20})
        # 실행이 끝나면 앱 설정 모델로 돌아온다
        self.assertNotEqual(analyzer.get_translation_model(), 'gemini-3.5-flash')

    def test_cron_run_is_skipped_when_the_switch_is_off(self):
        gmail_sync.set_cron_enabled(False)
        out = StringIO()
        with patch('api.management.commands.sync_gmail.sync_gmail') as mocked:
            call_command('sync_gmail', '--cron', stdout=out)
        mocked.assert_not_called()
        self.assertIn('꺼져 있습니다', out.getvalue())

    def test_manual_run_ignores_the_switch(self):
        # 웹 버튼/수동 실행은 스위치와 무관하게 동작한다
        gmail_sync.set_cron_enabled(False)
        with patch('api.management.commands.sync_gmail.sync_gmail',
                   return_value=dict(fetched=0, cases_created=0, emails_added=0,
                                     ignored=0, no_vendor=0, skipped=0, errors=0)) as mocked:
            call_command('sync_gmail', stdout=StringIO())
        mocked.assert_called_once()

    def test_concurrent_run_exits_quietly(self):
        # cron과 웹 버튼이 겹쳐도 cron 메일이 날아가지 않도록 정상 종료
        out = StringIO()
        with patch('api.management.commands.sync_gmail.sync_gmail',
                   side_effect=SyncInProgress('이미 진행 중입니다.')):
            call_command('sync_gmail', stdout=out)
        self.assertIn('이미 진행 중', out.getvalue())


class BackfillTranslationTests(TestCase):
    """backfill_translations — 번역 누락 메일만 채우고 케이스 필드는 불변."""

    def setUp(self):
        self.case = make_case(vendor='A10', summary='번역 백필 케이스',
                              status='Pending', action_steps='수동 편집 보존 확인')
        self.missing = make_email(self.case, 'Untranslated mail', message_id='bt-miss')
        self.missing.body_original = 'Hello, please check the device.'
        self.missing.save()
        self.translated = make_email(self.case, 'Translated mail', message_id='bt-done')
        self.translated.body_original = 'Already translated.'
        self.translated.subject_ko = '이미 번역된 메일'
        self.translated.body_ko = '이미 번역되어 있습니다.'
        self.translated.save()

    def run_command(self, *args):
        from io import StringIO
        from django.core.management import call_command
        out = StringIO()
        # 재시도 백오프(15초)로 테스트가 느려지지 않게 sleep 무력화
        with patch('api.management.commands.backfill_translations.time.sleep'):
            call_command('backfill_translations', '--sleep', '0', *args,
                         stdout=out, stderr=StringIO())
        return out.getvalue()

    def test_fills_only_missing_and_keeps_case_fields(self):
        with patch('api.management.commands.backfill_translations.analyze_email',
                   return_value={'subject_ko': '미번역 메일', 'body_ko': '장비를 확인해 주세요.'}) as mocked:
            self.run_command()
        self.missing.refresh_from_db()
        self.translated.refresh_from_db()
        self.case.refresh_from_db()
        self.assertEqual(self.missing.subject_ko, '미번역 메일')
        self.assertEqual(self.missing.body_ko, '장비를 확인해 주세요.')
        # 이미 번역된 메일은 호출조차 하지 않는다
        self.assertEqual(mocked.call_count, 1)
        self.assertEqual(self.translated.body_ko, '이미 번역되어 있습니다.')
        # 케이스 필드는 재분석하지 않으므로 그대로
        self.assertEqual(self.case.status, 'Pending')
        self.assertEqual(self.case.action_steps, '수동 편집 보존 확인')

    def test_failure_leaves_email_untouched_and_reports(self):
        with patch('api.management.commands.backfill_translations.analyze_email',
                   return_value=None) as mocked:
            out = self.run_command('--retries', '1')
        self.missing.refresh_from_db()
        self.assertEqual(self.missing.body_ko, '')
        self.assertEqual(mocked.call_count, 2)  # 원 시도 + 재시도 1회
        self.assertIn('실패 1건', out)


class ReferenceApiTests(TestCase):
    """Documents API — 목록/다운로드/업로드/임베딩/삭제와 역할별 권한 경계."""

    def setUp(self):
        from django.conf import settings as dj_settings
        from .permissions import set_user_role

        self.tmp = tempfile.TemporaryDirectory()
        self.addCleanup(self.tmp.cleanup)
        self.docs_root = Path(self.tmp.name)
        docs_override = override_settings(REFERENCE_DOCS_DIR=self.docs_root)
        docs_override.enable()
        self.addCleanup(docs_override.disable)
        self.embedding_model = dj_settings.EMBEDDING_MODEL

        self.viewer = User.objects.create_user('doc-viewer', password='x')
        set_user_role(self.viewer, 'viewer')
        self.engineer = User.objects.create_user('doc-eng', password='x')
        set_user_role(self.engineer, 'engineer')
        self.admin = User.objects.create_user('doc-admin', password='x')
        set_user_role(self.admin, 'admin')

    def _make_pdf(self, relative):
        path = self.docs_root / relative
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_bytes(b'%PDF-1.4 test content')
        return path

    def _upload(self, **extra):
        from django.core.files.uploadedfile import SimpleUploadedFile
        data = {
            'file': SimpleUploadedFile('New Guide.pdf', b'%PDF-1.4 uploaded',
                                       content_type='application/pdf'),
            'vendor': 'A10',
            'doc_type': 'config',
        }
        data.update(extra)
        return self.client.post('/api/references/upload/', data)

    def test_list_merges_files_and_embedding_state(self):
        from .models import ReferenceChunk, ReferenceDocument
        self._make_pdf('A10/config/embedded.pdf')
        self._make_pdf('Arista/release/raw.pdf')
        doc = ReferenceDocument.objects.create(
            vendor='A10', filename='A10/config/embedded.pdf', doc_type='config',
            title='ACOS Guide', sha256='x', page_count=10, chunk_count=3,
            embedding_model=self.embedding_model)
        # 임베딩 여부는 실제 청크 행으로 판정되므로 청크까지 만들어야 한다
        ReferenceChunk.objects.bulk_create([
            ReferenceChunk(document=doc, seq=i, page_start=i + 1, page_end=i + 1,
                           text=f'본문 {i}', embedding=b'\x00' * 16,
                           embedding_model=self.embedding_model)
            for i in range(3)
        ])

        self.client.force_login(self.viewer)  # 목록은 전 역할 공개
        data = self.client.get('/api/references/').json()
        self.assertEqual(len(data['items']), 2)
        by_name = {i['filename']: i for i in data['items']}
        self.assertTrue(by_name['A10/config/embedded.pdf']['embedded'])
        self.assertEqual(by_name['A10/config/embedded.pdf']['chunk_count'], 3)
        self.assertEqual(by_name['A10/config/embedded.pdf']['title'], 'ACOS Guide')
        self.assertFalse(by_name['Arista/release/raw.pdf']['embedded'])
        self.assertEqual(data['pending'], 1)
        self.assertTrue(data['auto_embed'])  # 미설정 기본값은 켜짐

    def test_download_inline_and_attachment(self):
        self._make_pdf('A10/config/guide.pdf')
        self.client.force_login(self.viewer)
        response = self.client.get('/api/references/file/',
                                   {'path': 'A10/config/guide.pdf'})
        self.assertEqual(response.status_code, 200)
        self.assertIn('inline', response['Content-Disposition'])
        response.close()
        response = self.client.get('/api/references/file/',
                                   {'path': 'A10/config/guide.pdf', 'dl': '1'})
        self.assertIn('attachment', response['Content-Disposition'])
        response.close()

    def test_download_blocks_traversal_and_unknown(self):
        self._make_pdf('A10/config/guide.pdf')
        self.client.force_login(self.viewer)
        for bad in ('../manage.py', '/etc/passwd', 'A10/../../manage.py', 'nope.pdf'):
            self.assertEqual(
                self.client.get('/api/references/file/', {'path': bad}).status_code,
                404, msg=bad)

    def test_upload_engineer_only_saves_and_conflicts(self):
        self.client.force_login(self.viewer)
        self.assertEqual(self._upload().status_code, 403)

        self.client.force_login(self.engineer)
        with patch('api.services.references.ingest_file',
                   return_value='created') as ingest:
            response = self._upload()
        self.assertEqual(response.status_code, 201, response.content)
        self.assertTrue((self.docs_root / 'A10/config/New Guide.pdf').is_file())
        self.assertTrue(response.json()['embedded'])  # 자동 임베딩 기본 켜짐
        ingest.assert_called_once()

        # 같은 이름 재업로드는 409, overwrite 지정 시 통과
        with patch('api.services.references.ingest_file', return_value='updated'):
            self.assertEqual(self._upload().status_code, 409)
            self.assertEqual(self._upload(overwrite='true').status_code, 201)

    def test_upload_skips_embedding_when_disabled(self):
        from .services import references
        references.set_auto_embed_enabled(False)
        self.client.force_login(self.engineer)
        with patch('api.services.references.ingest_file') as ingest:
            response = self._upload()
        self.assertEqual(response.status_code, 201)
        self.assertFalse(response.json()['embedded'])
        ingest.assert_not_called()

    def test_upload_validation(self):
        from django.core.files.uploadedfile import SimpleUploadedFile
        self.client.force_login(self.engineer)
        bad_ext = SimpleUploadedFile('run.exe', b'MZ')
        self.assertEqual(self._upload(file=bad_ext).status_code, 400)
        self.assertEqual(self._upload(vendor='NoSuchVendor').status_code, 400)
        with patch('api.services.references.ingest_file', return_value='created'):
            self.assertEqual(self._upload(doc_type='../evil').status_code, 400)

    def test_embed_admin_only(self):
        self._make_pdf('A10/config/guide.pdf')
        self.client.force_login(self.engineer)
        self.assertEqual(self.client.post('/api/references/embed/', {},
                                          content_type='application/json').status_code,
                         403)

        self.client.force_login(self.admin)
        with patch('api.services.references.ingest_file',
                   return_value='created') as ingest:
            data = self.client.post('/api/references/embed/', {},
                                    content_type='application/json').json()
        self.assertEqual(data['created'], 1)
        ingest.assert_called_once()

        # 단일 파일은 강제 재임베딩
        with patch('api.services.references.ingest_file',
                   return_value='updated') as ingest:
            response = self.client.post(
                '/api/references/embed/', {'path': 'A10/config/guide.pdf'},
                content_type='application/json')
        self.assertEqual(response.json()['outcome'], 'updated')
        self.assertTrue(ingest.call_args.kwargs['force'])

        response = self.client.post('/api/references/embed/', {'path': 'ghost.pdf'},
                                    content_type='application/json')
        self.assertEqual(response.status_code, 404)

    def test_delete_admin_only_removes_file_and_db(self):
        from .models import ReferenceDocument
        self._make_pdf('A10/config/guide.pdf')
        ReferenceDocument.objects.create(
            vendor='A10', filename='A10/config/guide.pdf', doc_type='config',
            sha256='x', chunk_count=1, embedding_model=self.embedding_model)

        self.client.force_login(self.engineer)
        self.assertEqual(
            self.client.delete('/api/references/file/?path=A10/config/guide.pdf')
            .status_code, 403)

        self.client.force_login(self.admin)
        response = self.client.delete(
            '/api/references/file/?path=A10/config/guide.pdf')
        self.assertEqual(response.status_code, 200)
        self.assertFalse((self.docs_root / 'A10/config/guide.pdf').exists())
        self.assertFalse(ReferenceDocument.objects.filter(
            filename='A10/config/guide.pdf').exists())

    def test_auto_embed_toggle_admin_only(self):
        self.client.force_login(self.viewer)
        self.assertTrue(self.client.get(
            '/api/settings/reference-auto-embed/').json()['enabled'])
        self.assertEqual(
            self.client.put('/api/settings/reference-auto-embed/',
                            json.dumps({'enabled': False}),
                            content_type='application/json').status_code, 403)

        self.client.force_login(self.admin)
        response = self.client.put('/api/settings/reference-auto-embed/',
                                   json.dumps({'enabled': False}),
                                   content_type='application/json')
        self.assertFalse(response.json()['enabled'])
        # boolean이 아니면 거부
        self.assertEqual(
            self.client.put('/api/settings/reference-auto-embed/',
                            json.dumps({'enabled': 'yes'}),
                            content_type='application/json').status_code, 400)


class ReferenceIngestIntegrityTests(TestCase):
    """인제스트 중간 실패가 DB를 망가뜨리지 않는지 — 2026-08-10 실장애 회귀 테스트.

    실장애: ClearPass 가이드 추출 텍스트에 NUL(0x00)이 섞여 bulk_create가
    DataError로 실패했는데, 트랜잭션이 없어 문서 행만 남고 청크는 0개가 됐다.
    목록은 chunk_count 필드를 믿어 "임베딩됨"으로 거짓 표시됐다.
    """

    def setUp(self):
        from django.conf import settings as dj_settings
        self.model = dj_settings.EMBEDDING_MODEL
        self.tmp = tempfile.TemporaryDirectory()
        self.addCleanup(self.tmp.cleanup)
        self.path = Path(self.tmp.name) / 'guide.pdf'
        self.path.write_bytes(b'%PDF-1.4 fake')

    def _ingest(self, pages):
        """extract_pages/embed_texts를 대역으로 두고 실제 저장 경로만 실행."""
        import numpy as np
        from .services import references
        vectors = np.zeros((max(len(pages), 1), 4), dtype=np.float32)
        with patch.object(references, 'extract_pages', return_value=pages), \
             patch.object(references, 'embed_texts', return_value=vectors):
            return references.ingest_file('A10', 'config', 'A10/config/guide.pdf',
                                          self.path, force=True)

    def test_nul_bytes_are_stripped_before_save(self):
        from .models import ReferenceChunk
        from .services import references
        self.assertEqual(references.sanitize('a\x00b'), 'ab')
        # 추출 단계에서 걸러지므로 저장되는 청크에는 NUL이 없다
        self._ingest([(1, references.sanitize('로그\x00 확인'))])
        text = ReferenceChunk.objects.get().text
        self.assertNotIn('\x00', text)
        self.assertEqual(text, '로그 확인')

    def test_failed_ingest_leaves_no_ghost_document(self):
        """저장 실패 시 문서 행도 남지 않아야 한다 (청크 0개 + 임베딩됨 표시 방지)."""
        from .models import ReferenceChunk, ReferenceDocument
        with patch.object(ReferenceChunk.objects, 'bulk_create',
                          side_effect=Exception('DataError: NUL byte')):
            with self.assertRaises(Exception):
                self._ingest([(1, '본문')])
        self.assertFalse(ReferenceDocument.objects.exists())
        self.assertFalse(ReferenceChunk.objects.exists())

    def test_failed_reingest_preserves_existing_chunks(self):
        """재임베딩이 깨져도 기존 청크가 삭제되면 안 된다 (삭제→삽입 순서 함정)."""
        from .models import ReferenceChunk, ReferenceDocument
        self._ingest([(1, '원본 내용')])
        self.assertEqual(ReferenceChunk.objects.count(), 1)

        with patch.object(ReferenceChunk.objects, 'bulk_create',
                          side_effect=Exception('DataError: NUL byte')):
            with self.assertRaises(Exception):
                self._ingest([(1, '새 내용')])

        self.assertEqual(ReferenceDocument.objects.count(), 1)
        chunks = ReferenceChunk.objects.all()
        self.assertEqual(len(chunks), 1)
        self.assertEqual(chunks[0].text, '원본 내용')  # 롤백되어 원본 유지

    def test_list_reports_ghost_document_as_not_embedded(self):
        """chunk_count 필드가 남아 있어도 실제 청크가 없으면 미임베딩으로 보여야 한다."""
        from django.contrib.auth.models import User
        from .models import ReferenceDocument
        from .permissions import set_user_role

        docs_root = Path(self.tmp.name) / 'refs'
        (docs_root / 'A10' / 'config').mkdir(parents=True)
        (docs_root / 'A10' / 'config' / 'ghost.pdf').write_bytes(b'%PDF-1.4')
        ReferenceDocument.objects.create(
            vendor='A10', filename='A10/config/ghost.pdf', doc_type='config',
            sha256='x', page_count=253, chunk_count=97,  # 필드만 남은 유령 행
            embedding_model=self.model)

        admin = User.objects.create_user('ghost-admin', password='x')
        set_user_role(admin, 'admin')
        self.client.force_login(admin)
        with override_settings(REFERENCE_DOCS_DIR=docs_root):
            data = self.client.get('/api/references/').json()

        item = data['items'][0]
        self.assertFalse(item['embedded'])
        self.assertEqual(item['chunk_count'], 0)
        self.assertEqual(data['pending'], 1)


class HelpAgentProgressCallbackTests(TestCase):
    """on_event 콜백이 진행 상황을 순서대로 흘리는지 (스트리밍의 원천).

    콜백을 넘기지 않으면 기존과 완전히 동일하게 동작해야 한다 — 이 보장 덕분에
    chat()을 호출하는 기존 테스트 21곳이 그대로 통과한다.
    """

    def _fake_client(self, tool_names):
        """도구를 tool_names 순서로 한 번씩 호출한 뒤 최종 답변을 내는 가짜 클라이언트."""
        calls = []

        def create(**kwargs):
            step = len(calls)
            calls.append(kwargs)
            if step < len(tool_names):
                block = SimpleNamespace(type='tool_use', id=f'tu{step}',
                                        name=tool_names[step], input={'query': '비밀검색어'})
                return SimpleNamespace(stop_reason='tool_use', content=[block])
            return SimpleNamespace(
                stop_reason='end_turn',
                content=[SimpleNamespace(type='text', text='최종 답변')])

        return SimpleNamespace(messages=SimpleNamespace(create=create))

    def test_events_follow_tool_order(self):
        events = []
        client = self._fake_client(['search_knowledge', 'search_references'])
        with patch.object(help_agent, '_execute_tool', return_value=('결과', False)):
            reply, trace, _, _ = help_agent._run_agent(
                client, 'search', [{'role': 'user', 'content': '질문'}],
                on_event=lambda kind, payload: events.append((kind, payload)))

        self.assertEqual(reply, '최종 답변')
        kinds = [k for k, _ in events]
        self.assertEqual(kinds, ['step', 'tool', 'step', 'tool', 'step'])
        tools = [p['name'] for k, p in events if k == 'tool']
        self.assertEqual(tools, ['search_knowledge', 'search_references'])
        # 반복 횟수를 함께 실어 UI가 진행감을 표시할 수 있어야 한다
        self.assertEqual(events[1][1]['iteration'], 1)
        self.assertEqual(events[3][1]['iteration'], 2)

    def test_tool_input_is_never_leaked(self):
        """도구 입력에는 고객사명·시리얼이 섞일 수 있어 이벤트로 내보내면 안 된다."""
        events = []
        client = self._fake_client(['web_search'])
        with patch.object(help_agent, '_execute_tool', return_value=('결과', False)):
            help_agent._run_agent(client, 'search', [{'role': 'user', 'content': 'q'}],
                                  on_event=lambda kind, payload: events.append((kind, payload)))
        payloads = json.dumps([p for _, p in events], ensure_ascii=False)
        self.assertNotIn('비밀검색어', payloads)
        self.assertNotIn('input', payloads)

    def test_without_callback_behaves_as_before(self):
        client = self._fake_client(['search_cases'])
        with patch.object(help_agent, '_execute_tool', return_value=('결과', False)):
            reply, trace, _, _ = help_agent._run_agent(
                client, 'search', [{'role': 'user', 'content': '질문'}])
        self.assertEqual(reply, '최종 답변')
        self.assertEqual([t['name'] for t in trace], ['search_cases'])


class HelpAgentChatStreamTests(TestCase):
    """POST /api/help-agent/chat/stream/ — SSE 스트리밍 응답."""

    FAKE_RESULT = {
        'reply': '스트리밍 답변', 'tool_calls': [{'name': 'search_cases', 'input': {}}],
        'model': 'claude-haiku-4-5', 'agent': 'search',
    }

    def setUp(self):
        from .permissions import set_user_role
        viewer = User.objects.create_user('s-viewer', password='pw123456')
        set_user_role(viewer, 'viewer')
        engineer = User.objects.create_user('s-eng', password='pw123456')
        set_user_role(engineer, 'engineer')

    def post(self, payload):
        return self.client.post('/api/help-agent/chat/stream/', payload,
                                content_type='application/json')

    @staticmethod
    def read_events(response):
        body = b''.join(response.streaming_content).decode()
        events = []
        for block in body.split('\n\n'):
            if not block.strip():
                continue
            kind = next(l[6:].strip() for l in block.split('\n') if l.startswith('event:'))
            data = next(l[5:].strip() for l in block.split('\n') if l.startswith('data:'))
            events.append((kind, json.loads(data)))
        return events

    def test_viewer_is_blocked(self):
        self.client.force_login(User.objects.get(username='s-viewer'))
        res = self.post({'messages': [{'role': 'user', 'content': '안녕'}]})
        self.assertEqual(res.status_code, 403)

    def test_validation_happens_before_stream_starts(self):
        """검증 실패는 스트림이 아니라 상태 코드로 와야 한다 (이벤트 오류로 감추면 안 됨)."""
        self.client.force_login(User.objects.get(username='s-eng'))
        self.assertEqual(self.post({'messages': []}).status_code, 400)
        res = self.post({'messages': [{'role': 'user', 'content': '안녕'}],
                         'session_id': 99999})
        self.assertEqual(res.status_code, 404)

    def test_progress_events_then_done(self):
        self.client.force_login(User.objects.get(username='s-eng'))

        def fake_chat(messages, on_event=None):
            on_event('step', {'phase': 'triage'})
            on_event('tool', {'name': 'search_cases', 'iteration': 1, 'max': 6})
            return dict(self.FAKE_RESULT)

        # 저장은 워커 스레드가 별도 DB 커넥션으로 수행하므로 TestCase의 트랜잭션에서는
        # 보이지 않는다. 저장 로직 자체는 비스트리밍 뷰 테스트가 이미 덮으므로,
        # 여기서는 "스트림이 끝나기 전에 저장을 호출했는가"만 검증한다.
        with patch('api.views.help_agent.chat', side_effect=fake_chat), \
             patch('api.views.HelpAgentChatStreamView._save_turns',
                   return_value=42) as save, \
             patch('api.views.log_event') as logged:
            res = self.post({'messages': [{'role': 'user', 'content': 'VRRP 사례'}]})
            self.assertEqual(res.status_code, 200)
            self.assertEqual(res['Content-Type'], 'text/event-stream')
            events = self.read_events(res)

        kinds = [k for k, _ in events]
        self.assertEqual(kinds, ['step', 'tool', 'done'])
        done = events[-1][1]
        # done 페이로드는 비스트리밍 응답과 동일해야 한다 (프론트 렌더링 재사용)
        self.assertEqual(done['reply'], '스트리밍 답변')
        self.assertEqual(done['agent'], 'search')
        self.assertEqual(done['session_id'], 42)
        save.assert_called_once()
        self.assertEqual(logged.call_args.args[1], 'agent_chat')

    def test_runtime_failure_becomes_error_event(self):
        import httpx
        self.client.force_login(User.objects.get(username='s-eng'))
        error = anthropic.RateLimitError(
            'rate limited',
            response=httpx.Response(429, request=httpx.Request('POST', 'https://x')),
            body=None)
        with patch('api.views.help_agent.chat', side_effect=error), \
             patch('api.views.HelpAgentChatStreamView._save_turns') as save:
            res = self.post({'messages': [{'role': 'user', 'content': '안녕'}]})
            self.assertEqual(res.status_code, 200)  # 헤더는 이미 나갔다
            events = self.read_events(res)

        self.assertEqual([k for k, _ in events], ['error'])
        self.assertEqual(events[0][1]['code'], 'rate_limit')
        self.assertIn('사용량 한도', events[0][1]['message'])
        save.assert_not_called()  # 실패한 응답은 저장하지 않는다


class HelpAgentDocumentAccessTests(TestCase):
    """검색 에이전트가 벤더 문서를 조회할 수 있는지 — 2026-08-11 실사용 신고 회귀.

    신고: "사내 문서 검색을 통해 확인해줘"에 검색 에이전트가 "저는 케이스 DB만
    검색할 수 있다, SharePoint는 IT에 문의하라"고 거절. 실제로는 시스템에 벤더
    문서 2만여 청크가 적재돼 있었고, 원인은 검색 에이전트에 search_references
    도구가 없었던 것 + 프롬프트의 거절 지시가 SCOPE_GUARD의 핸드오프 규칙과
    충돌한 것이었다.
    """

    def test_search_agent_can_search_documents(self):
        tools = {t['name'] for t in help_agent._agent_configs()['search']['tools']}
        self.assertIn('search_references', tools)
        # 케이스 조회 능력은 그대로 유지
        self.assertIn('search_cases', tools)
        self.assertIn('search_knowledge', tools)

    def test_search_prompt_does_not_instruct_refusal(self):
        prompt = help_agent._agent_configs()['search']['system']
        # 거절 대신 핸드오프가 원칙 (SCOPE_GUARD와 충돌하던 문장 제거됨)
        self.assertNotIn('일반 기술 지원은 범위 밖', prompt)
        self.assertIn('HANDOFF:tech', prompt)
        # 외부 문서 시스템으로 떠넘기지 않도록 명시
        self.assertIn('SharePoint', prompt)
        self.assertIn('search_references', prompt)

    def test_missing_document_returns_explicit_notice(self):
        """문서에 없는 버그 번호는 빈 결과가 아니라 '못 찾았다'는 안내를 돌려준다."""
        with patch('api.services.references.search', return_value=[]):
            payload = json.loads(help_agent._search_references(query='ACOS-104904'))
        self.assertEqual(payload['count'], 0)
        self.assertIn('찾지 못했습니다', payload['notice'])


class CaseKnowledgeExtractTests(TestCase):
    """POST /api/cases/<id>/knowledge/ — 케이스 상세의 '지식으로 저장' 버튼.

    벤더 확답은 케이스 종결 전에 나오는 일이 많아, Resolved만 스캔하는 자동
    동기화로는 그동안 정보가 묶인다 (2026-08-11 C-1118 사례). 상태와 무관하게
    엔지니어가 직접 뽑을 수 있어야 한다.
    """

    def setUp(self):
        from .permissions import set_user_role
        self.case = make_case(status='Pending', summary='진행 중 케이스')
        viewer = User.objects.create_user('k-viewer', password='x')
        set_user_role(viewer, 'viewer')
        engineer = User.objects.create_user('k-eng', password='x')
        set_user_role(engineer, 'engineer')
        self.engineer = engineer

    def post(self):
        return self.client.post(f'/api/cases/{self.case.id}/knowledge/')

    def test_viewer_is_blocked(self):
        self.client.force_login(User.objects.get(username='k-viewer'))
        self.assertEqual(self.post().status_code, 403)

    def test_extracts_from_ongoing_case(self):
        from .models import KnowledgeItem
        self.client.force_login(self.engineer)
        with patch('api.services.knowledge.generate_structured_with_model', return_value=(MODEL, {
            'has_knowledge': True, 'title': 'ACOS-104904 VRRP 타이머 버그',
            'problem': 'VRRP advertisement 타이머 부정확',
            'root_cause': 'ACOS 6.0.6-SP1~6.0.8 버그',
            'resolution': '6.0.9 이상으로 업그레이드',
            'device_model': 'TH1040-F', 'software_version': '6.0.8',
        })), patch('api.services.knowledge.enrich_with_references'):
            res = self.post()

        self.assertEqual(res.status_code, 201)
        self.assertEqual(res.json()['outcome'], 'created')
        item = KnowledgeItem.objects.get()
        self.assertEqual(item.status, 'draft')  # 확정은 사람이 한다
        self.assertEqual(item.case, self.case)

    def test_ongoing_case_is_not_marked_checked(self):
        """진행 중 케이스를 '검토 완료'로 찍으면 해결된 뒤 자동 동기화가 영영 건너뛴다."""
        self.client.force_login(self.engineer)
        with patch('api.services.knowledge.generate_structured_with_model',
                   return_value=(MODEL, {'has_knowledge': False, 'resolution': ''})):
            res = self.post()

        self.assertEqual(res.status_code, 400)
        self.assertEqual(res.json()['outcome'], 'no_knowledge')
        self.case.refresh_from_db()
        self.assertIsNone(self.case.knowledge_checked_at)

    def test_resolved_case_is_marked_checked(self):
        """해결된 케이스는 기존대로 검토 완료 표시 — 자동 동기화가 재검토하지 않게."""
        self.case.status = 'Resolved'
        self.case.save(update_fields=['status'])
        self.client.force_login(self.engineer)
        with patch('api.services.knowledge.generate_structured_with_model',
                   return_value=(MODEL, {'has_knowledge': False, 'resolution': ''})):
            self.post()
        self.case.refresh_from_db()
        self.assertIsNotNone(self.case.knowledge_checked_at)

    def test_second_press_returns_existing_item(self):
        self.client.force_login(self.engineer)
        payload = {
            'has_knowledge': True, 'title': '제목', 'problem': '문제',
            'root_cause': '원인', 'resolution': '해결',
            'device_model': '', 'software_version': '',
        }
        with patch('api.services.knowledge.generate_structured_with_model',
                   return_value=(MODEL, payload)), \
             patch('api.services.knowledge.enrich_with_references'):
            first = self.post()
            second = self.post()

        self.assertEqual(first.status_code, 201)
        self.assertEqual(second.status_code, 200)
        self.assertEqual(second.json()['outcome'], 'exists')
        self.assertEqual(
            first.json()['item']['knowledge_id'], second.json()['item']['knowledge_id'])

    def test_missing_case_returns_404(self):
        self.client.force_login(self.engineer)
        self.assertEqual(
            self.client.post('/api/cases/999999/knowledge/').status_code, 404)
